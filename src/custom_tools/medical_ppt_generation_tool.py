"""
医疗病例PPT生成工具
使用python-pptx库在本地生成医疗病例PPT文件
"""

from typing import Any, Type, Optional, List, Dict
from pydantic import BaseModel, Field
from crewai.tools import BaseTool
import os
import logging
import uuid
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

from src.custom_tools.medical_ppt_template import get_slide_schema, get_template_by_id
from src.utils.qiniu_client import QiniuClient

logger = logging.getLogger(__name__)

# 医疗主题配色方案 - 蓝白配色
MEDICAL_COLORS = {
    'primary': RGBColor(41, 128, 185),      # 医疗蓝 - 主色
    'primary_dark': RGBColor(31, 97, 141),  # 深蓝 - 强调
    'primary_light': RGBColor(52, 152, 219),# 亮蓝 - 标题栏
    'accent': RGBColor(66, 139, 202),       # 中蓝 - 强调色（原绿色改为蓝色）
    'accent_light': RGBColor(100, 181, 246),# 浅蓝 - 装饰
    'warning': RGBColor(230, 126, 34),      # 橙色 - 警告
    'background': RGBColor(245, 248, 250),  # 极浅蓝灰 - 背景
    'text_primary': RGBColor(44, 62, 80),   # 深灰 - 主文本
    'text_secondary': RGBColor(127, 140, 141),  # 中灰 - 次要文本
    'white': RGBColor(255, 255, 255),       # 纯白色
    'table_header': RGBColor(41, 128, 185), # 表头颜色 - 主蓝色
    'table_alt_row': RGBColor(240, 248, 255)  # 表格交替行 - 淡蓝色
}


def truncate_text(text: str, max_length: int, suffix: str = "...") -> str:
    """
    截断文本到指定长度
    
    Args:
        text: 原始文本
        max_length: 最大长度
        suffix: 截断后缀
        
    Returns:
        截断后的文本
    """
    if not text:
        return ""
    text = str(text)
    if len(text) <= max_length:
        return text
    return text[:max_length - len(suffix)] + suffix


class MedicalPPTGenerationToolSchema(BaseModel):
    """医疗PPT生成工具输入Schema"""
    slides_data: List[Dict[str, Any]] = Field(
        ..., description="""幻灯片数据列表，每个幻灯片必须包含layout和content字段:
    - layout: 必须使用医疗PPT模板中定义的幻灯片id (例如: 'medical:patient-info-slide')
    - content: 必须严格遵循对应layout的json_schema规范"""
    )
    patient_name: str = Field(
        default="患者", description="患者姓名，用于演示文稿标题"
    )
    template_id: str = Field(
        default="medical", description="模板ID (目前支持: 'medical')"
    )
    output_dir: str = Field(
        default="./output/presentation", description="输出目录"
    )
    upload_to_qiniu: bool = Field(
        default=True, description="是否上传生成的PPT到七牛云"
    )
    session_id: str = Field(
        default="", description="会话ID，用于组织文件"
    )


class MedicalPPTGenerationTool(BaseTool):
    name: str = "Generate Medical Case PowerPoint Presentation"
    description: str = (
        "本地生成医疗病例PPT的工具。"
        "接收结构化的幻灯片数据，使用python-pptx库在本地创建PPTX文件。"
        "支持病例信息、诊断、检查、治疗等多种医疗专用页面模板。"
    )
    args_schema: Type[BaseModel] = MedicalPPTGenerationToolSchema
    result_as_answer: bool = True  # 工具返回结果直接作为最终答案

    def _run(self, **kwargs: Any) -> Any:
        """执行PPT生成"""
        slides_data = kwargs.get("slides_data")
        patient_name = kwargs.get("patient_name", "患者")
        template_id = kwargs.get("template_id", "medical")
        output_dir = kwargs.get("output_dir", "./output/presentation")
        upload_to_qiniu = kwargs.get("upload_to_qiniu", True)
        session_id = kwargs.get("session_id", "")

        # 验证slides_data
        if not slides_data or not isinstance(slides_data, list):
            return {"success": False, "error": "slides_data必须是非空列表"}

        # 验证每个幻灯片的必需字段
        for i, slide in enumerate(slides_data):
            if not isinstance(slide, dict):
                return {"success": False, "error": f"幻灯片 {i} 必须是字典类型"}
            if 'layout' not in slide:
                return {"success": False, "error": f"幻灯片 {i} 缺少必需字段 'layout'"}
            if 'content' not in slide:
                return {"success": False, "error": f"幻灯片 {i} 缺少必需字段 'content'"}

        print(f"开始生成医疗病例PPT: 患者 {patient_name}")
        print(f"模板: {template_id}, 幻灯片数量: {len(slides_data)}")

        try:
            # 创建演示文稿
            result = self._create_presentation(
                patient_name=patient_name,
                template_id=template_id,
                slides_data=slides_data,
                output_dir=output_dir,
                session_id=session_id
            )

            if not result:
                return {"success": False, "error": "创建演示文稿失败"}

            # 上传到七牛云（如果启用）
            qiniu_url = None
            file_uuid = None
            file_key = None
            if upload_to_qiniu and result.get('local_path'):
                try:
                    qiniu_result = self._upload_to_qiniu(
                        local_path=result['local_path'],
                        session_id=session_id,
                        export_as="pptx"
                    )
                    if qiniu_result:
                        qiniu_url = qiniu_result.get('url')
                        file_uuid = qiniu_result.get('file_uuid')
                        file_key = qiniu_result.get('file_key')
                        print(f"成功上传到七牛云: {qiniu_url}")
                        print(f"文件UUID: {file_uuid}")
                        print(f"文件Key: {file_key}")
                except Exception as e:
                    logger.error(f"上传到七牛云失败: {str(e)}")
                    print(f"警告: 上传到七牛云失败: {str(e)}")

            print(f"PPT生成完成")
            return {
                "success": True,
                "local_path": result.get('local_path'),
                "filename": result.get('filename'),
                "qiniu_url": qiniu_url,
                "file_uuid": file_uuid,
                "file_key": file_key
            }

        except Exception as e:
            logger.error(f"PPT生成失败: {str(e)}", exc_info=True)
            print(f"PPT生成过程中出错: {str(e)}")
            return {"success": False, "error": str(e)}

    def _create_presentation(
        self,
        patient_name: str,
        template_id: str,
        slides_data: List[Dict[str, Any]],
        output_dir: str,
        session_id: str = ""
    ) -> Optional[Dict[str, Any]]:
        """
        创建演示文稿
        
        Args:
            patient_name: 患者姓名
            template_id: 模板ID
            slides_data: 幻灯片数据列表
            output_dir: 输出目录
            session_id: 会话ID
            
        Returns:
            包含本地路径和文件名的字典
        """
        try:
            # 创建演示文稿对象
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            print(f"创建演示文稿: 患者MDT病例整理 - {patient_name}")
            print(f"  模板: {template_id}")
            print(f"  幻灯片数量: {len(slides_data)}")

            # 添加标题页
            self._add_title_slide(prs, patient_name)

            # 遍历添加每个幻灯片
            for i, slide_data in enumerate(slides_data):
                layout_id = slide_data.get('layout')
                content = slide_data.get('content')
                
                print(f"  添加幻灯片 {i+1}: {layout_id}")
                
                # 根据layout类型添加相应的幻灯片
                if layout_id == "medical:patient-info-slide":
                    self._add_patient_info_slide(prs, content)
                elif layout_id == "medical:diagnosis-list-slide":
                    self._add_diagnosis_list_slide(prs, content)
                elif layout_id == "medical:lab-test-table-slide":
                    self._add_lab_test_table_slide(prs, content)
                elif layout_id == "medical:metric-trends-slide":
                    self._add_metric_trends_slide(prs, content)
                elif layout_id == "medical:imaging-table-slide":
                    self._add_imaging_table_slide(prs, content)
                elif layout_id == "medical:imaging-comparison-slide":
                    self._add_imaging_comparison_slide(prs, content)
                elif layout_id == "medical:treatment-history-slide":
                    self._add_treatment_history_slide(prs, content)
                elif layout_id == "medical:timeline-slide":
                    self._add_timeline_slide(prs, content)
                else:
                    print(f"    警告: 未知的layout类型: {layout_id}")

            # 创建输出目录
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            if session_id:
                output_path = os.path.join(output_dir, session_id)
            else:
                output_path = output_dir
            os.makedirs(output_path, exist_ok=True)

            # 保存文件
            filename = f"medical_case_{patient_name}_{timestamp}.pptx"
            local_path = os.path.join(output_path, filename)
            prs.save(local_path)

            print(f"演示文稿已保存: {local_path}")
            
            return {
                'local_path': local_path,
                'filename': filename
            }

        except Exception as e:
            logger.error(f"创建演示文稿失败: {str(e)}", exc_info=True)
            print(f"创建演示文稿失败: {str(e)}")
            return None

    def _add_standard_title_bar(self, slide, title_text: str):
        """为幻灯片添加统一的标题栏（优化版）
        
        Args:
            slide: 幻灯片对象
            title_text: 标题文本
        """
        # 截断标题文本
        title_text = truncate_text(title_text, 20)
        
        # 标题背景 - 使用渐变效果的背景条
        title_bg = slide.shapes.add_shape(1, Inches(0.4), Inches(0.25), Inches(9.2), Inches(0.8))
        title_bg_fill = title_bg.fill
        title_bg_fill.solid()
        title_bg_fill.fore_color.rgb = MEDICAL_COLORS['primary_light']
        
        # 添加边框阴影效果
        title_bg.line.fill.background()
        title_bg.shadow.inherit = False
        
        # 左侧装饰条 - 使用深蓝色
        deco_bar = slide.shapes.add_shape(1, Inches(0.4), Inches(0.25), Inches(0.15), Inches(0.8))
        deco_fill = deco_bar.fill
        deco_fill.solid()
        deco_fill.fore_color.rgb = MEDICAL_COLORS['primary_dark']
        deco_bar.line.fill.background()
        
        # 标题文字
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.3), Inches(8.5), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = MEDICAL_COLORS['white']
        title_frame.paragraphs[0].font.name = "Microsoft YaHei"
        title_frame.vertical_anchor = 1  # 居中对齐
    
    def _add_white_background(self, slide):
        """为幻灯片添加白色背景
        
        Args:
            slide: 幻灯片对象
        """
        bg_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
        bg_fill = bg_shape.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = RGBColor(255, 255, 255)
        bg_shape.line.fill.background()

    def _add_title_slide(self, prs: Presentation, patient_name: str):
        """添加标题页 - 优化版本"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 使用空白布局
        
        # 截断患者姓名
        patient_name = truncate_text(patient_name, 10)
        
        # 添加背景色块 - 白色背景
        bg_shape = slide.shapes.add_shape(
            1,  # 矩形
            Inches(0), Inches(0),
            Inches(10), Inches(7.5)
        )
        bg_fill = bg_shape.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = MEDICAL_COLORS['white']
        bg_shape.line.fill.background()
        
        # 添加顶部装饰区域 - 使用渐变色
        top_bar = slide.shapes.add_shape(
            1,  # 矩形
            Inches(0), Inches(0),
            Inches(10), Inches(2)
        )
        top_fill = top_bar.fill
        top_fill.solid()
        top_fill.fore_color.rgb = MEDICAL_COLORS['primary']
        top_bar.line.fill.background()
        
        # 添加底部装饰条 - 使用浅蓝色
        bottom_bar = slide.shapes.add_shape(
            1,  # 矩形
            Inches(0), Inches(7.3),
            Inches(10), Inches(0.2)
        )
        bottom_fill = bottom_bar.fill
        bottom_fill.solid()
        bottom_fill.fore_color.rgb = MEDICAL_COLORS['accent_light']
        bottom_bar.line.fill.background()
        
        # 主标题
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "患者MDT病例整理"
        title_frame.paragraphs[0].font.size = Pt(48)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = MEDICAL_COLORS['primary']
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.name = "Microsoft YaHei"
        
        # 患者信息卡片 - 使用白色背景和蓝色边框
        card_bg = slide.shapes.add_shape(
            1,  # 矩形
            Inches(2.5), Inches(4.2),
            Inches(5), Inches(1.2)
        )
        card_fill = card_bg.fill
        card_fill.solid()
        card_fill.fore_color.rgb = MEDICAL_COLORS['white']
        card_bg.line.color.rgb = MEDICAL_COLORS['primary']
        card_bg.line.width = Pt(3)
        
        # 患者信息文字
        info_box = slide.shapes.add_textbox(Inches(2.5), Inches(4.3), Inches(5), Inches(1))
        info_frame = info_box.text_frame
        info_frame.text = f"患者：{patient_name}"
        info_frame.paragraphs[0].font.size = Pt(28)
        info_frame.paragraphs[0].font.bold = True
        info_frame.paragraphs[0].font.color.rgb = MEDICAL_COLORS['text_primary']
        info_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        info_frame.paragraphs[0].font.name = "Microsoft YaHei"
        info_frame.vertical_anchor = 1
        
        # 日期
        date_box = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(8), Inches(0.5))
        date_frame = date_box.text_frame
        date_frame.text = f"报告日期：{datetime.now().strftime('%Y年%m月%d日')}"
        date_frame.paragraphs[0].font.size = Pt(18)
        date_frame.paragraphs[0].font.color.rgb = MEDICAL_COLORS['text_secondary']
        date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        date_frame.paragraphs[0].font.name = "Microsoft YaHei"

    def _add_patient_info_slide(self, prs: Presentation, content: Dict[str, Any]):
        """添加病例信息页 - 优化版"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 使用空白布局
        
        # 添加白色背景和标题
        self._add_white_background(slide)
        self._add_standard_title_bar(slide, content.get('title', '病例信息'))
        
        # 内容区域 - 使用卡片式布局
        y_offset = 1.3
        left_margin = 0.8
        label_width = 2.2
        content_width = 6.5
        row_height = 0.85
        
        info_items = [
            ("姓名", truncate_text(content.get('patient_name', ''), 10)),
            ("性别", truncate_text(content.get('gender', ''), 10)),
            ("年龄", truncate_text(content.get('age', ''), 20)),
            ("既往史", truncate_text(content.get('medical_history', ''), 200)),
            ("个人史", truncate_text(content.get('personal_history', ''), 200)),
            ("会诊目的", truncate_text(content.get('consultation_purpose', ''), 200))
        ]
        
        for i, (label, value) in enumerate(info_items):
            # 交替背景色 - 使用淡蓝色
            if i % 2 == 0:
                bg_box = slide.shapes.add_shape(
                    1,
                    Inches(left_margin - 0.2),
                    Inches(y_offset - 0.05),
                    Inches(9),
                    Inches(row_height)
                )
                bg_fill = bg_box.fill
                bg_fill.solid()
                bg_fill.fore_color.rgb = MEDICAL_COLORS['table_alt_row']
                bg_box.line.fill.background()
            
            # 标签
            label_box = slide.shapes.add_textbox(
                Inches(left_margin), 
                Inches(y_offset), 
                Inches(label_width), 
                Inches(row_height - 0.1)
            )
            label_frame = label_box.text_frame
            label_frame.text = label + "："
            label_frame.paragraphs[0].font.size = Pt(16)
            label_frame.paragraphs[0].font.bold = True
            label_frame.paragraphs[0].font.color.rgb = MEDICAL_COLORS['primary']
            label_frame.paragraphs[0].font.name = "Microsoft YaHei"
            label_frame.word_wrap = True
            label_frame.vertical_anchor = 1
            
            # 内容
            content_box = slide.shapes.add_textbox(
                Inches(left_margin + label_width), 
                Inches(y_offset), 
                Inches(content_width), 
                Inches(row_height - 0.1)
            )
            content_frame = content_box.text_frame
            content_frame.text = str(value) if value else "未填写"
            content_frame.paragraphs[0].font.size = Pt(14)
            content_frame.paragraphs[0].font.color.rgb = MEDICAL_COLORS['text_primary']
            content_frame.paragraphs[0].font.name = "Microsoft YaHei"
            content_frame.word_wrap = True
            content_frame.vertical_anchor = 1
            
            y_offset += row_height

    def _add_diagnosis_list_slide(self, prs: Presentation, content: Dict[str, Any]):
        """添加诊断信息页 - 优化版"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 使用空白布局
        
        # 添加白色背景和标题
        self._add_white_background(slide)
        self._add_standard_title_bar(slide, content.get('title', '诊断信息'))
        
        # 诊断列表
        diagnoses = content.get('diagnoses', [])
        y_offset = 1.4
        left_margin = 1.0
        
        for i, diagnosis in enumerate(diagnoses):
            index = diagnosis.get('index', i + 1)
            diag_text = truncate_text(diagnosis.get('diagnosis', ''), 100)
            diag_time = truncate_text(diagnosis.get('diagnosis_time', ''), 30)
            
            # 背景卡片 - 使用白色背景和淡蓝色边框
            card_bg = slide.shapes.add_shape(
                1,
                Inches(left_margin - 0.2),
                Inches(y_offset - 0.05),
                Inches(8.5),
                Inches(0.65)
            )
            card_fill = card_bg.fill
            card_fill.solid()
            card_fill.fore_color.rgb = MEDICAL_COLORS['white']
            card_bg.line.color.rgb = MEDICAL_COLORS['accent_light']
            card_bg.line.width = Pt(1.5)
            
            # 序号圆圈
            circle = slide.shapes.add_shape(
                1,
                Inches(left_margin),
                Inches(y_offset + 0.05),
                Inches(0.4),
                Inches(0.4)
            )
            circle_fill = circle.fill
            circle_fill.solid()
            circle_fill.fore_color.rgb = MEDICAL_COLORS['primary']
            circle.line.fill.background()
            
            # 序号文字
            num_box = slide.shapes.add_textbox(
                Inches(left_margin),
                Inches(y_offset + 0.05),
                Inches(0.4),
                Inches(0.4)
            )
            num_frame = num_box.text_frame
            num_frame.text = str(index)
            num_frame.paragraphs[0].font.size = Pt(14)
            num_frame.paragraphs[0].font.bold = True
            num_frame.paragraphs[0].font.color.rgb = MEDICAL_COLORS['white']
            num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            num_frame.paragraphs[0].font.name = "Microsoft YaHei"
            num_frame.vertical_anchor = 1
            
            # 诊断文字
            diag_box = slide.shapes.add_textbox(
                Inches(left_margin + 0.6),
                Inches(y_offset),
                Inches(7.2),
                Inches(0.55)
            )
            diag_frame = diag_box.text_frame
            if diag_time:
                diag_frame.text = f"{diag_text}  [{diag_time}]"
            else:
                diag_frame.text = diag_text
            
            diag_frame.paragraphs[0].font.size = Pt(16)
            diag_frame.paragraphs[0].font.color.rgb = MEDICAL_COLORS['text_primary']
            diag_frame.paragraphs[0].font.name = "Microsoft YaHei"
            diag_frame.word_wrap = True
            diag_frame.vertical_anchor = 1
            
            y_offset += 0.75

    def _add_lab_test_table_slide(self, prs: Presentation, content: Dict[str, Any]):
        """添加实验室检查表格页 - 优化版"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 使用空白布局
        
        # 添加白色背景和标题
        self._add_white_background(slide)
        self._add_standard_title_bar(slide, content.get('title', '实验室检查'))
        
        # 表格数据
        table_data = content.get('table_data', {})
        headers = table_data.get('headers', ['检查项目', '开单时间', '项目名称', '结果', '报告时间'])
        rows = table_data.get('rows', [])
        
        if rows:
            # 添加表格
            cols = len(headers)
            table_rows = len(rows) + 1  # +1 for header
            
            # 计算表格位置和大小
            left = Inches(0.5)
            top = Inches(1.4)
            width = Inches(9)
            height = Inches(5.8)
            
            table = slide.shapes.add_table(table_rows, cols, left, top, width, height).table
            
            # 设置表头样式
            for col_idx, header in enumerate(headers):
                header_text = truncate_text(header, 30)
                cell = table.cell(0, col_idx)
                cell.text = header_text
                cell.text_frame.paragraphs[0].font.size = Pt(12)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
                cell.fill.solid()
                cell.fill.fore_color.rgb = MEDICAL_COLORS['table_header']
                cell.text_frame.paragraphs[0].font.color.rgb = MEDICAL_COLORS['white']
                cell.vertical_anchor = 1
            
            # 填充数据行
            for row_idx, row_data in enumerate(rows, 1):
                for col_idx, cell_data in enumerate(row_data):
                    if col_idx < cols:  # 确保不超出列数
                        # 截断单元格文本
                        cell_text = truncate_text(str(cell_data), 100)
                        cell = table.cell(row_idx, col_idx)
                        cell.text = cell_text
                        cell.text_frame.paragraphs[0].font.size = Pt(10)
                        cell.text_frame.paragraphs[0].font.color.rgb = MEDICAL_COLORS['text_primary']
                        cell.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
                        cell.text_frame.word_wrap = True
                        cell.vertical_anchor = 1
                        
                        # 交替行颜色
                        if row_idx % 2 == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = MEDICAL_COLORS['table_alt_row']
                        else:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = MEDICAL_COLORS['white']

    def _add_metric_trends_slide(self, prs: Presentation, content: Dict[str, Any]):
        """添加关键指标趋势图页"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 使用空白布局
        
        # 添加白色背景和标题
        self._add_white_background(slide)
        self._add_standard_title_bar(slide, content.get('title', '关键指标变动趋势'))
        
        # 获取指标数据
        metrics = content.get('metrics', [])
        
        if not metrics:
            # 如果没有数据，显示提示信息
            text_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
            text_frame = text_box.text_frame
            text_frame.text = "暂无指标数据"
            text_frame.paragraphs[0].font.size = Pt(24)
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            return
        
        # 为每个指标创建折线图
        y_start = 1.2
        chart_height = (6.5 - y_start) / len(metrics) - 0.2
        
        for i, metric in enumerate(metrics):
            metric_name = metric.get('name', f'指标{i+1}')
            unit = metric.get('unit', '')
            normal_range = metric.get('normal_range', '')
            data_points = metric.get('data_points', [])
            
            if not data_points:
                continue
            
            # 准备图表数据
            chart_data = CategoryChartData()
            chart_data.categories = [dp.get('date', '') for dp in data_points]
            chart_data.add_series(f'{metric_name}', [dp.get('value', 0) for dp in data_points])
            
            # 添加折线图
            y_pos = y_start + i * (chart_height + 0.2)
            x, y, cx, cy = Inches(0.5), Inches(y_pos), Inches(8.5), Inches(chart_height)
            
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
            ).chart
            
            # 设置图表标题
            chart.has_title = True
            chart_title = f"{metric_name}"
            if unit:
                chart_title += f" ({unit})"
            if normal_range:
                chart_title += f" 正常范围: {normal_range}"
            chart.chart_title.text_frame.text = chart_title
            chart.chart_title.text_frame.paragraphs[0].font.size = Pt(12)

    def _add_imaging_table_slide(self, prs: Presentation, content: Dict[str, Any]):
        """添加影像学检查表格页 - 优化版"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 使用空白布局
        
        # 添加白色背景和标题
        self._add_white_background(slide)
        self._add_standard_title_bar(slide, content.get('title', '影像学检查'))
        
        # 表格数据
        table_data = content.get('table_data', {})
        headers = table_data.get('headers', ['检查项目', '报告内容', '检查时间'])
        rows = table_data.get('rows', [])
        
        if rows:
            cols = len(headers)
            table_rows = len(rows) + 1
            
            left = Inches(0.5)
            top = Inches(1.4)
            width = Inches(9)
            height = Inches(4.8)
            
            table = slide.shapes.add_table(table_rows, cols, left, top, width, height).table
            
            # 设置表头样式
            for col_idx, header in enumerate(headers):
                header_text = truncate_text(header, 30)
                cell = table.cell(0, col_idx)
                cell.text = header_text
                cell.text_frame.paragraphs[0].font.size = Pt(12)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
                cell.fill.solid()
                cell.fill.fore_color.rgb = MEDICAL_COLORS['table_header']
                cell.text_frame.paragraphs[0].font.color.rgb = MEDICAL_COLORS['white']
                cell.vertical_anchor = 1
            
            # 数据行
            for row_idx, row_data in enumerate(rows, 1):
                for col_idx, cell_data in enumerate(row_data):
                    if col_idx < cols:
                        # 截断单元格文本
                        cell_text = truncate_text(str(cell_data), 100)
                        cell = table.cell(row_idx, col_idx)
                        cell.text = cell_text
                        cell.text_frame.paragraphs[0].font.size = Pt(10)
                        cell.text_frame.paragraphs[0].font.color.rgb = MEDICAL_COLORS['text_primary']
                        cell.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
                        cell.text_frame.word_wrap = True
                        cell.vertical_anchor = 1
                        
                        # 交替行颜色
                        if row_idx % 2 == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = MEDICAL_COLORS['table_alt_row']
                        else:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = MEDICAL_COLORS['white']
        
        # 备注
        note = content.get('note', '')
        if note:
            note_text = truncate_text(note, 200)
            note_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.4), Inches(8.8), Inches(0.8))
            note_frame = note_box.text_frame
            note_frame.text = f"备注：{note_text}"
            note_frame.paragraphs[0].font.size = Pt(12)
            note_frame.paragraphs[0].font.color.rgb = MEDICAL_COLORS['text_secondary']
            note_frame.paragraphs[0].font.name = "Microsoft YaHei"
            note_frame.paragraphs[0].font.italic = True
            note_frame.word_wrap = True

    def _add_imaging_comparison_slide(self, prs: Presentation, content: Dict[str, Any]):
        """添加影像学对比页 - 优化版"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 使用空白布局
        
        # 添加白色背景和标题
        self._add_white_background(slide)
        self._add_standard_title_bar(slide, content.get('title', '影像学检查'))
        
        # 副标题
        subtitle = content.get('subtitle', '')
        if subtitle:
            subtitle_text = truncate_text(subtitle, 50)
            subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(8.8), Inches(0.3))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle_text
            subtitle_frame.paragraphs[0].font.size = Pt(12)
            subtitle_frame.paragraphs[0].font.color.rgb = MEDICAL_COLORS['text_secondary']
            subtitle_frame.paragraphs[0].font.name = "Microsoft YaHei"
            subtitle_frame.paragraphs[0].font.italic = True
        
        # 图片列表
        images = content.get('images', [])
        
        if not images:
            return
        
        # 根据图片数量确定布局
        num_images = len(images)
        
        if num_images == 1:
            # 单张图片，居中显示
            self._add_single_image(slide, images[0], Inches(2), Inches(1.8), Inches(6), Inches(4.5))
        elif num_images == 2:
            # 两张图片，左右排列
            self._add_single_image(slide, images[0], Inches(0.5), Inches(1.8), Inches(4.5), Inches(4.5))
            self._add_single_image(slide, images[1], Inches(5.2), Inches(1.8), Inches(4.5), Inches(4.5))
        elif num_images == 3:
            # 三张图片，一张在上，两张在下
            self._add_single_image(slide, images[0], Inches(3), Inches(1.5), Inches(4), Inches(2.5))
            self._add_single_image(slide, images[1], Inches(0.5), Inches(4.2), Inches(4.5), Inches(2.5))
            self._add_single_image(slide, images[2], Inches(5.2), Inches(4.2), Inches(4.5), Inches(2.5))
        else:  # 4张或更多
            # 四张图片，2x2网格
            positions = [
                (Inches(0.5), Inches(1.8), Inches(4.5), Inches(2.5)),
                (Inches(5.2), Inches(1.8), Inches(4.5), Inches(2.5)),
                (Inches(0.5), Inches(4.5), Inches(4.5), Inches(2.5)),
                (Inches(5.2), Inches(4.5), Inches(4.5), Inches(2.5))
            ]
            for i, img in enumerate(images[:4]):  # 最多显示4张
                left, top, width, height = positions[i]
                self._add_single_image(slide, img, left, top, width, height)

    def _add_single_image(self, slide, image_data: Dict, left, top, width, height):
        """添加单张图片到幻灯片 - 优化版"""
        image_path = image_data.get('image_path', '')
        caption = image_data.get('caption', '')
        
        try:
            if image_path and os.path.exists(image_path):
                # 添加图片
                pic = slide.shapes.add_picture(image_path, left, top, width=width, height=height)
                
                # 添加图片说明
                if caption:
                    caption_text = truncate_text(caption, 50)
                    caption_box = slide.shapes.add_textbox(
                        left, 
                        top + height + Inches(0.05), 
                        width, 
                        Inches(0.3)
                    )
                    caption_frame = caption_box.text_frame
                    caption_frame.text = caption_text
                    caption_frame.paragraphs[0].font.size = Pt(9)
                    caption_frame.paragraphs[0].font.color.rgb = MEDICAL_COLORS['text_secondary']
                    caption_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    caption_frame.paragraphs[0].font.name = "Microsoft YaHei"
            else:
                # 如果图片不存在，显示占位符
                placeholder_box = slide.shapes.add_textbox(left, top, width, height)
                placeholder_frame = placeholder_box.text_frame
                placeholder_frame.text = f"图片不存在\n{image_path}" if image_path else "未提供图片"
                placeholder_frame.paragraphs[0].font.size = Pt(14)
                placeholder_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                
                # 添加边框
                line = placeholder_box.line
                line.color.rgb = RGBColor(150, 150, 150)
                line.width = Pt(1)
                
        except Exception as e:
            logger.error(f"添加图片失败: {str(e)}")
            print(f"添加图片失败: {str(e)}")

    def _add_treatment_history_slide(self, prs: Presentation, content: Dict[str, Any]):
        """添加治疗历史页 - 优化版"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 使用空白布局
        
        # 添加白色背景和标题
        self._add_white_background(slide)
        self._add_standard_title_bar(slide, content.get('title', '治疗历史'))
        
        # 表格数据
        table_data = content.get('table_data', {})
        headers = table_data.get('headers', ['治疗名称', '治疗记录', '开始时间', '结束时间'])
        rows = table_data.get('rows', [])
        
        if rows:
            cols = len(headers)
            table_rows = len(rows) + 1
            
            left = Inches(0.5)
            top = Inches(1.4)
            width = Inches(9)
            height = Inches(5.8)
            
            table = slide.shapes.add_table(table_rows, cols, left, top, width, height).table
            
            # 设置表头样式
            for col_idx, header in enumerate(headers):
                header_text = truncate_text(header, 30)
                cell = table.cell(0, col_idx)
                cell.text = header_text
                cell.text_frame.paragraphs[0].font.size = Pt(12)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
                cell.fill.solid()
                cell.fill.fore_color.rgb = MEDICAL_COLORS['table_header']
                cell.text_frame.paragraphs[0].font.color.rgb = MEDICAL_COLORS['white']
                cell.vertical_anchor = 1
            
            # 数据行
            for row_idx, row_data in enumerate(rows, 1):
                for col_idx, cell_data in enumerate(row_data):
                    if col_idx < cols:
                        # 截断单元格文本
                        cell_text = truncate_text(str(cell_data), 100)
                        cell = table.cell(row_idx, col_idx)
                        cell.text = cell_text
                        cell.text_frame.paragraphs[0].font.size = Pt(10)
                        cell.text_frame.paragraphs[0].font.color.rgb = MEDICAL_COLORS['text_primary']
                        cell.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
                        cell.text_frame.word_wrap = True
                        cell.vertical_anchor = 1
                        
                        # 交替行颜色
                        if row_idx % 2 == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = MEDICAL_COLORS['table_alt_row']
                        else:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = MEDICAL_COLORS['white']

    def _add_timeline_slide(self, prs: Presentation, content: Dict[str, Any]):
        """添加患者疾病旅程时间轴页 - 优化版"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 使用空白布局
        
        # 添加白色背景和标题
        self._add_white_background(slide)
        self._add_standard_title_bar(slide, content.get('title', '患者疾病旅程(汇总)'))
        
        # 如果有时间轴图片，直接显示图片
        timeline_image_path = content.get('timeline_image_path', '')
        if timeline_image_path and os.path.exists(timeline_image_path):
            try:
                slide.shapes.add_picture(
                    timeline_image_path, 
                    Inches(0.5), 
                    Inches(1.3), 
                    width=Inches(9),
                    height=Inches(5.8)
                )
                return
            except Exception as e:
                logger.error(f"添加时间轴图片失败: {str(e)}")
        
        # 如果没有图片，使用文本方式显示事件
        events = content.get('events', [])
        if events:
            y_offset = 1.5
            left_margin = 0.8
            
            for i, event in enumerate(events):
                date = truncate_text(event.get('date', ''), 30)
                event_text = truncate_text(event.get('event', ''), 100)
                
                # 时间线连接线
                if i > 0:
                    line = slide.shapes.add_shape(
                        1,
                        Inches(left_margin + 0.15),
                        Inches(y_offset - 0.4),
                        Inches(0.05),
                        Inches(0.35)
                    )
                    line_fill = line.fill
                    line_fill.solid()
                    line_fill.fore_color.rgb = MEDICAL_COLORS['primary_light']
                    line.line.fill.background()
                
                # 时间点标记 - 使用中蓝色
                marker = slide.shapes.add_shape(
                    1,
                    Inches(left_margin),
                    Inches(y_offset),
                    Inches(0.35),
                    Inches(0.35)
                )
                marker_fill = marker.fill
                marker_fill.solid()
                marker_fill.fore_color.rgb = MEDICAL_COLORS['accent']
                marker.line.fill.background()
                
                # 日期
                date_box = slide.shapes.add_textbox(
                    Inches(left_margin + 0.5),
                    Inches(y_offset),
                    Inches(2.5),
                    Inches(0.35)
                )
                date_frame = date_box.text_frame
                date_frame.text = date
                date_frame.paragraphs[0].font.size = Pt(13)
                date_frame.paragraphs[0].font.bold = True
                date_frame.paragraphs[0].font.color.rgb = MEDICAL_COLORS['primary']
                date_frame.paragraphs[0].font.name = "Microsoft YaHei"
                date_frame.vertical_anchor = 1
                
                # 事件内容
                event_box = slide.shapes.add_textbox(
                    Inches(left_margin + 3.2),
                    Inches(y_offset),
                    Inches(5.5),
                    Inches(0.35)
                )
                event_frame = event_box.text_frame
                event_frame.text = event_text
                event_frame.paragraphs[0].font.size = Pt(13)
                event_frame.paragraphs[0].font.color.rgb = MEDICAL_COLORS['text_primary']
                event_frame.paragraphs[0].font.name = "Microsoft YaHei"
                event_frame.word_wrap = True
                event_frame.vertical_anchor = 1
                
                y_offset += 0.65
        else:
            # 显示提示文本
            text_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(6), Inches(1))
            text_frame = text_box.text_frame
            text_frame.text = "暂无时间轴数据"
            text_frame.paragraphs[0].font.size = Pt(20)
            text_frame.paragraphs[0].font.color.rgb = MEDICAL_COLORS['text_secondary']
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            text_frame.paragraphs[0].font.name = "Microsoft YaHei"

    def _upload_to_qiniu(self, local_path: str, session_id: str, export_as: str) -> Optional[Dict[str, str]]:
        """
        上传文件到七牛云并返回公开URL和file_uuid
        
        Args:
            local_path: 本地文件路径
            session_id: 会话ID
            export_as: 导出格式(如'pptx')
            
        Returns:
            包含url和file_uuid的字典，失败返回None
        """
        try:
            # 创建七牛云客户端
            qiniu_client = QiniuClient()

            # 生成UUID作为文件名
            file_uuid = str(uuid.uuid4())
            
            # 生成文件key，包含session_id和UUID
            if session_id:
                file_key = f"presentations/{session_id}/{file_uuid}.{export_as}"
            else:
                file_key = f"presentations/{file_uuid}.{export_as}"

            # 上传文件
            ret, info = qiniu_client.upload_file(local_path, key=file_key)

            if info.status_code == 200:
                # 获取公开下载URL
                download_url = qiniu_client.download_file(file_key, is_private=False)

                logger.info(f"文件成功上传到七牛云: file_uuid={file_uuid}, file_key={file_key}")

                return {
                    "url": download_url,
                    "file_uuid": file_uuid,
                    "file_key": file_key  # 保存完整的file_key路径，前端用于生成下载URL
                }
            else:
                logger.error(f"七牛云上传失败，状态码: {info.status_code}")
                return None

        except Exception as e:
            logger.error(f"上传到七牛云时出错: {str(e)}")
            return None

