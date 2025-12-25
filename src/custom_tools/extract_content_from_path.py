# src/custom_tools/extract_content_from_path.py

from typing import Any, Optional, Type
from pydantic import BaseModel, Field
from crewai.tools import BaseTool
import os
import re
import zipfile
import requests
import docx
import html2text
from goose3 import Goose
from goose3.text import StopWordsChinese
from bs4 import BeautifulSoup
import trafilatura
import chardet
import random
# import pypdf
import pypdf
import pdfplumber
import fitz  # PyMuPDF
import io
import logging
from dotenv import load_dotenv
import json
import base64  # 新增导入
from openai import OpenAI  # 新增导入
import shutil  # 新增导入，用于清理临时文件夹
import time  # 新增导入，用于时间统计
import concurrent.futures  # 🚨 新增：用于并发处理
from functools import partial  # 🚨 新增：用于创建偏函数
import uuid  # 新增导入，用于生成UUID
from pathlib import Path  # 新增导入，用于处理文件路径



load_dotenv()

# Set up logger
logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)
handler = logging.StreamHandler()
formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
handler.setFormatter(formatter)
logger.addHandler(handler)

# 尝试导入HEIC转换相关库
try:
    from PIL import Image
    from pillow_heif import register_heif_opener
    register_heif_opener()
    HEIC_CONVERSION_AVAILABLE = True
    logger.info("HEIC转换功能可用")
except ImportError:
    HEIC_CONVERSION_AVAILABLE = False
    logger.warning("HEIC转换功能不可用，需要安装 pillow-heif 库")
    
try:
    from pptx import Presentation  # 新增导入，用于PPT文件处理
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available. PPT file processing will be skipped.")
    
# 初始化OpenAI客户端用于图片处理
try:
    # 从环境变量获取配置
    api_keys_str = os.getenv("MULTIMODAL_API_KEY")
    base_url = os.getenv("MULTIMODAL_BASE_URL")
    multimodal_model_name_str = os.getenv("MULTIMODAL_MODEL_NAME")

    # 解析多个 API keys（逗号分隔）
    if api_keys_str:
        multimodal_api_keys = [key.strip() for key in api_keys_str.split(',') if key.strip()]
    else:
        multimodal_api_keys = []

    # 解析多个模型名称（逗号分隔）
    if multimodal_model_name_str:
        multimodal_model_names = [name.strip() for name in multimodal_model_name_str.split(',') if name.strip()]
    else:
        multimodal_model_names = []

    if not multimodal_api_keys or not base_url:
        logger.warning("MULTIMODAL_API_KEY or MULTIMODAL_BASE_URL not found in environment variables")
        openai_client = None
    else:
        # 使用第一个 API key 初始化默认客户端（实际使用时会动态选择）
        openai_client = OpenAI(
            api_key=multimodal_api_keys[0],
            base_url=base_url
        )
        logger.info(f"OpenAI client initialized with {len(multimodal_api_keys)} API key(s) and {len(multimodal_model_names)} model(s): {multimodal_model_names}")
except Exception as e:
    logger.warning(f"Failed to initialize OpenAI client for image processing: {e}")
    openai_client = None
    multimodal_api_keys = []
    multimodal_model_names = []

class ExtractContentFromPathsSchema(BaseModel):
    """Input for ExtractContentFromPathsTool."""
    paths: list[str] = Field(..., description="List of absolute paths or URLs to extract content from")

class ExtractContentFromPathsTool(BaseTool):
    name: str = "Extract content from multiple files path or URLs"
    description: str = (
        "A tool that can extract content from multiple files (PDF, DOCX, TXT, images) or URLs at once."
    )
    args_schema: Type[BaseModel] = ExtractContentFromPathsSchema

    def _run(
        self,
        paths: list[str],
    ) -> Any:
        single_extractor = ExtractContentFromPathTool()
        results = []
        
        for path in paths:
            try:
                result = single_extractor._run(path=path)
                if result:
                    if isinstance(result, list):
                        results.extend(result)
                    else:
                        results.append(result)
            except Exception as e:
                logger.info(f"Error processing path {path}: {e}")
                continue
                
        return json.dumps(results)
    
class ExtractContentFromPathSchema(BaseModel):
    """Input for ExtractContentFromPathTool."""
    path: str = Field(..., description="Absolute Path to the file or URL to extract content from")

class ExtractContentFromPathTool(BaseTool):
    name: str = "Extract content from file or URL"
    description: str = (
        "A tool that can extract content from various file types (PDF, DOCX, TXT, images) or URLs."
    )
    args_schema: Type[BaseModel] = ExtractContentFromPathSchema

    def _run(
        self,
        path: str,
    ) -> Any:
        # 过滤系统隐藏文件
        if os.path.isfile(path):
            filename = os.path.basename(path)
            if filename.startswith('._') or filename.startswith('.DS_Store'):
                logger.warning(f"跳过系统隐藏文件: {filename}")
                return {
                    'file_extension': 'hidden',
                    'file_name': filename,
                    'file_content': f"系统隐藏文件: {filename} (已跳过处理)",
                    'extraction_success': False,
                    'extraction_error': '系统隐藏文件，已跳过处理'
                    # 注意：不生成UUID，由上层统一管理
                }

        try:
            if path.startswith(('http://', 'https://')):
                result = self.read_url(path)
            elif os.path.isfile(path):
                file_extension = os.path.splitext(path)[1].lower()
                if file_extension == '.pdf':
                    # 根据环境变量选择PDF处理方式
                    pdf_extraction_mode = os.getenv("PDF_EXTRACTION_MODE", "default").lower()

                    if pdf_extraction_mode == "with_images":
                        logger.info(f"使用图片提取模式处理PDF: {os.path.basename(path)}")
                        result = self.read_pdf_with_images(path)
                    else:
                        logger.info(f"使用默认模式处理PDF: {os.path.basename(path)}")
                        result = self.read_pdf(path)
                elif file_extension == '.docx':
                    result = self.read_docx(path)
                elif file_extension in ['.pptx', '.ppt']:  # 新增PPT文件支持
                    result = self.read_ppt(path)
                elif file_extension == '.txt':
                    result = self.read_txt(path)
                elif file_extension == '.md':
                    result = self.read_md(path)
                elif file_extension == '.json':
                    result = self.read_json(path)
                elif file_extension == '.zip':
                    # 解压zip文件并递归处理所有文件
                    result = self.process_zip_file(path)
                elif file_extension in ['.png', '.jpg', '.jpeg', '.webp', '.heic', '.heif']:
                    result = self.read_image(path)
                else:
                    result = self.read_file(path)
            else:
                raise ValueError(f"Invalid path or unsupported file type: {path}")

            # 🚨 标记提取状态（不生成UUID，UUID由上层file_processing_manager统一管理）
            if isinstance(result, dict):
                # 标记提取成功（如果result是字典且没有error字段）
                if 'extraction_success' not in result:
                    # 判断是否提取成功：有file_content且内容不为空
                    has_content = result.get('file_content') and len(str(result.get('file_content', '')).strip()) > 0
                    result['extraction_success'] = has_content
                    if not has_content:
                        result['extraction_error'] = '提取内容为空'

                # 注意：不在这里生成file_uuid，由上层统一管理
            elif isinstance(result, list):
                # 对于返回列表的情况（如zip文件、PDF with images），标记提取状态
                for item in result:
                    if isinstance(item, dict):
                        # 标记提取成功
                        if 'extraction_success' not in item:
                            has_content = item.get('file_content') and len(str(item.get('file_content', '')).strip()) > 0
                            item['extraction_success'] = has_content
                            if not has_content:
                                item['extraction_error'] = '提取内容为空'

                        # 注意：这里保留UUID生成，因为zip/PDF中的子文件需要新的UUID
                        # 但是主文件的UUID应该保留原始值
                        if 'file_uuid' not in item:
                            item['file_uuid'] = str(uuid.uuid4())

            return result
        except Exception as e:
            filename = os.path.basename(path) if os.path.isfile(path) else path
            logger.error(f"文件提取失败: {filename}, 错误: {str(e)}")
            return {
                'file_extension': os.path.splitext(filename)[1].lower()[1:] if os.path.isfile(path) else 'unknown',
                'file_name': filename,
                'file_content': f"文件提取失败: {str(e)}",
                'extraction_success': False,
                'extraction_error': f"{type(e).__name__}: {str(e)}"
                # 注意：不生成UUID，由上层统一管理
            }

    def detect_image_format(self, path):
        """检测图片的实际格式，不依赖文件扩展名"""
        try:
            # 读取文件头来检测格式
            with open(path, 'rb') as f:
                header = f.read(32)  # 读取更多字节以支持更多格式
                
            # 检查常见的图片格式魔数
            if header.startswith(b'\x89PNG\r\n\x1a\n'):
                return 'png'
            elif header.startswith(b'\xff\xd8\xff'):
                return 'jpeg'
            elif header.startswith(b'RIFF') and b'WEBP' in header:
                return 'webp'
            elif header.startswith(b'GIF87a') or header.startswith(b'GIF89a'):
                return 'gif'
            elif header.startswith(b'BM'):
                return 'bmp'
            elif b'ftypheic' in header or b'ftypheix' in header:
                return 'heic'
            elif b'ftypmif1' in header or b'ftypmsf1' in header:
                return 'heif'
            elif header.startswith(b'\x00\x00\x01\x00') or header.startswith(b'\x00\x00\x02\x00'):
                return 'ico'
                
        except Exception as e:
            logger.debug(f"格式检测失败: {e}")
            
        return None

    def convert_heic_to_jpg(self, heic_path, output_path=None):
        """将HEIC格式转换为JPG格式"""
        if not HEIC_CONVERSION_AVAILABLE:
            logger.error("HEIC转换功能不可用，需要安装 pillow-heif 库")
            return None
            
        try:
            # 如果没有指定输出路径，则在同目录下生成
            if output_path is None:
                base_name = os.path.splitext(heic_path)[0]
                output_path = f"{base_name}_converted.jpg"
            
            # 使用PIL打开HEIC文件并转换为JPG
            with Image.open(heic_path) as img:
                # 如果图片有透明通道，转换为RGB模式
                if img.mode in ('RGBA', 'LA', 'P'):
                    # 创建白色背景
                    background = Image.new('RGB', img.size, (255, 255, 255))
                    if img.mode == 'P':
                        img = img.convert('RGBA')
                    background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else None)
                    img = background
                elif img.mode != 'RGB':
                    img = img.convert('RGB')
                
                # 保存为JPG格式，质量设置为90
                img.save(output_path, 'JPEG', quality=90, optimize=True)
                
            logger.info(f"HEIC转换成功: {os.path.basename(heic_path)} -> {os.path.basename(output_path)}")
            return output_path
            
        except Exception as e:
            logger.error(f"HEIC转换失败: {str(e)}")
            return None

    def read_image(self, path):
        """使用多模态模型提取图片内容描述"""
        try:
            filename = os.path.basename(path)
            file_extension = os.path.splitext(path)[1].lower()
            result = {
                'file_extension': file_extension[1:],
                'file_name': filename,
                'extraction_success': False,  # 默认失败，成功后会更新
                'extraction_error': None
            }
            
            # 过滤隐藏文件和系统文件
            if filename.startswith('._') or filename.startswith('.DS_Store'):
                logger.warning(f"跳过系统隐藏文件: {filename}")
                result['file_content'] = f"系统隐藏文件: {filename} (已跳过处理)"
                result['extraction_error'] = '系统隐藏文件，已跳过处理'
                return result
            
            # 检查文件存在性和大小
            if not os.path.exists(path):
                result['file_content'] = f"图片文件不存在: {filename}"
                return result
                
            file_size = os.path.getsize(path)
            if file_size == 0:
                logger.warning(f"文件大小为0: {filename}")
                result['file_content'] = f"图片文件: {filename} (文件大小为0)"
                return result
            elif file_size > 20 * 1024 * 1024:  # 大于20MB
                logger.warning(f"文件过大: {filename} ({file_size / 1024 / 1024:.1f}MB)")
                result['file_content'] = f"图片文件: {filename} (文件过大: {file_size / 1024 / 1024:.1f}MB)"
                return result
            
            # 检测实际图片格式
            actual_format = self.detect_image_format(path)
            original_path = path  # 保存原始路径
            converted_file = None  # 转换后的文件路径
            
            if actual_format:
                # 如果实际格式与扩展名不符，记录警告
                expected_format = file_extension[1:] if file_extension else 'unknown'
                if actual_format != expected_format and expected_format != 'jpg':  # jpg和jpeg视为相同
                    if not (actual_format == 'jpeg' and expected_format in ['jpg', 'jpeg']):
                        logger.warning(f"格式不匹配: {filename} 扩展名={expected_format}, 实际={actual_format}")
                
                # 如果是HEIC格式，先转换为JPG
                if actual_format in ['heic', 'heif']:
                    logger.info(f"检测到HEIC/HEIF格式，开始转换: {filename}")
                    import tempfile
                    # 在临时目录创建转换后的文件
                    temp_dir = tempfile.mkdtemp()
                    converted_filename = f"{os.path.splitext(filename)[0]}_converted.jpg"
                    converted_path = os.path.join(temp_dir, converted_filename)
                    
                    converted_file = self.convert_heic_to_jpg(path, converted_path)
                    if converted_file:
                        path = converted_file  # 使用转换后的文件路径
                        actual_format = 'jpeg'  # 更新格式
                        logger.info(f"HEIC转换成功，使用转换后的文件进行处理")
                    else:
                        logger.error(f"HEIC转换失败，尝试直接处理原文件")
                
                # 使用实际格式确定MIME类型
                format_to_mime = {
                    'png': 'image/png',
                    'jpeg': 'image/jpeg',
                    'jpg': 'image/jpeg',
                    'webp': 'image/webp',
                    'heic': 'image/jpeg',  # HEIC使用JPEG MIME类型
                    'heif': 'image/jpeg'   # HEIF使用JPEG MIME类型
                }
                mime_type = format_to_mime.get(actual_format, 'image/png')
            else:
                # 回退到基于扩展名的MIME类型
                mime_type_map = {
                    '.png': 'image/png',
                    '.jpg': 'image/jpeg',
                    '.jpeg': 'image/jpeg',
                    '.webp': 'image/webp',
                    '.heic': 'image/jpeg',
                    '.heif': 'image/jpeg'
                }
                mime_type = mime_type_map.get(file_extension, 'image/png')
            
            # 检查多模态客户端
            if not openai_client:
                logger.warning(f"多模态客户端不可用: {filename}")
                result['file_content'] = f"图片文件: {filename} (无法处理图片内容，多模态模型不可用)"
                return result
            
            logger.info(f"开始处理图片: {filename} (格式: {actual_format or 'unknown'})")
            
            # 读取图片并转换为base64
            with open(path, 'rb') as image_file:
                image_data = image_file.read()
                image_base64 = base64.b64encode(image_data).decode('utf-8')
            
            # 调用多模态模型（支持多API key和多模型容错）
            try:
                response_content = None
                last_error = None
                successful_model = None
                successful_api_key_index = None

                # 随机打乱 API keys 顺序
                shuffled_api_keys = multimodal_api_keys.copy()
                random.shuffle(shuffled_api_keys)

                for api_key_index, api_key in enumerate(shuffled_api_keys):
                    # 为每个 API key 创建新的客户端
                    try:
                        current_client = OpenAI(
                            api_key=api_key,
                            base_url=base_url
                        )
                    except Exception as client_error:
                        logger.warning(f"创建客户端失败 (API key {api_key_index + 1}/{len(shuffled_api_keys)}): {str(client_error)[:100]}")
                        continue

                    # 尝试该 API key 的所有模型
                    for model_index, model_name in enumerate(multimodal_model_names):
                        try:
                            logger.info(f"尝试 API key {api_key_index + 1}/{len(shuffled_api_keys)}, 模型 {model_index + 1}/{len(multimodal_model_names)}: {model_name}")

                            # 尝试使用JSON模式（部分模型不支持，需要容错处理）
                            request_params = {
                                "model": model_name,
                                "messages": [
                                    {
                                        "role": "user",
                                        "content": [
                                            {"type": "text", "text": """请分析这张图片并返回JSON格式的结果:
{
  "has_medical_image": true/false,
  "content": "图片内容描述",
  "image_bbox": {
    "x": 0.0-1.0,
    "y": 0.0-1.0,
    "width": 0.0-1.0,
    "height": 0.0-1.0
  }
}

【判断规则】:
1. has_medical_image: 判断图片是否包含医学影像(如CT、MRI、X光片、超声、病理切片、内镜图像等)
   - 如果图片中包含实际的医学影像图片(不是纯文字报告截图),返回true
   - 如果是纯文字报告单、检验单截图,返回false
   - 如果不是医学相关图片,返回false

2. content: 图片内容描述
   - 只返回医学资料的内容,不返回其他无关的内容
   - 如果和医学资料无关,则不返回内容
   - 如果和医学资料有关,请详细描述这张图片的内容,包括图片中的文字、日期、或者对象,比如检查检验单,CT影像等所有可见信息
   - 如果图片包含文档或表格,请尽可能准确地转录其中的文字内容
   - 对于医学表格检验检查报告单,不需要一些医生姓名相关的信息
   - 结果返回解析后的医学报告的markdown格式内容

3. image_bbox: 医学影像图片的边界框(仅当has_medical_image=true时填写)
   - 用于从混合图片(文字+图像)中裁剪出医学影像部分
   - bbox坐标使用归一化值(0.0-1.0),其中:
     * x: 图像中心点的水平位置(0=最左,1=最右)
     * y: 图像中心点的垂直位置(0=最上,1=最下)
     * width: 图像宽度占整张图片宽度的比例
     * height: 图像高度占整张图片高度的比例
   - 如果整张图片就是医学影像(没有文字报告部分),bbox可以是 {"x": 0.5, "y": 0.5, "width": 1.0, "height": 1.0}
   - 如果图片是文字报告中嵌入的医学影像,返回影像部分的准确bbox

直接返回JSON,不要输出其他回复语气词"""},
                                        {
                                            "type": "image_url",
                                            "image_url": {
                                                "url": f"data:{mime_type};base64,{image_base64}"
                                            }
                                        }
                                    ]
                                }
                            ],
                            "max_tokens": 8190
                        }

                            # 仅对支持JSON模式的模型添加response_format参数
                            # GLM-4.5V等部分模型不支持该参数，会返回400错误
                            if "Qwen" in model_name or "gpt-" in model_name.lower():
                                request_params["response_format"] = {"type": "json_object"}

                            response = current_client.chat.completions.create(**request_params)

                            response_content = response.choices[0].message.content
                            successful_model = model_name
                            successful_api_key_index = api_key_index
                            logger.info(f"调用成功 - API key {api_key_index + 1}, 模型: {model_name}")
                            break  # 成功则跳出模型循环

                        except Exception as api_error:
                            last_error = api_error
                            logger.warning(f"调用失败 - API key {api_key_index + 1}, 模型 {model_name}: {str(api_error)[:200]}")

                            # 如果不是最后一个模型，继续尝试下一个模型
                            if model_index < len(multimodal_model_names) - 1:
                                logger.info(f"继续尝试该 API key 的下一个模型...")
                                continue

                    # 如果当前 API key 的某个模型成功，跳出 API key 循环
                    if response_content:
                        break

                    # 如果当前 API key 的所有模型都失败，尝试下一个 API key
                    if api_key_index < len(shuffled_api_keys) - 1:
                        logger.info(f"该 API key 所有模型均失败，尝试下一个 API key...")
                    else:
                        logger.error(f"所有 API key 和模型组合均失败")
                        if last_error:
                            raise last_error

                # 处理响应内容
                if response_content:
                    # 解析JSON响应
                    try:
                        import json as json_lib
                        response_data = json_lib.loads(response_content.strip())

                        has_medical_image = response_data.get('has_medical_image', False)
                        image_description = response_data.get('content', '')
                        image_bbox = response_data.get('image_bbox')  # 简化：只保留一个bbox

                        # 清理特殊标记
                        if image_description:
                            image_description = image_description.strip()
                            if image_description.startswith("\n<|begin_of_box|>"):
                                image_description = image_description[len("\n<|begin_of_box|>"):].strip()
                            elif image_description.startswith("<|begin_of_box|>"):
                                image_description = image_description[len("<|begin_of_box|>"):].strip()
                            if image_description.endswith("<|end_of_box|>"):
                                image_description = image_description[:-len("<|end_of_box|>")].strip()

                        if image_description.strip():
                            log_msg = f"成功提取图片内容: {filename} ({len(image_description)} 字符), 模型: {successful_model}, 是否包含医学影像: {has_medical_image}"
                            if image_bbox:
                                log_msg += f", 边界框: x={image_bbox.get('x', 0):.2f}, y={image_bbox.get('y', 0):.2f}"
                            logger.info(log_msg)

                            result['file_content'] = image_description
                            result['has_medical_image'] = has_medical_image
                            result['model_used'] = successful_model
                            result['extraction_success'] = True  # 标记提取成功

                            # 存储边界框用于裁剪
                            if has_medical_image and image_bbox:
                                result['image_bbox'] = image_bbox
                                logger.info(f"医学影像边界框: {image_bbox}")

                                # 裁剪医学影像并保存
                                try:
                                    from PIL import Image as PILImage
                                    img = PILImage.open(path)
                                    img_width, img_height = img.size

                                    # 转换归一化坐标为像素坐标
                                    center_x = image_bbox.get('x', 0.5) * img_width
                                    center_y = image_bbox.get('y', 0.5) * img_height
                                    bbox_width = image_bbox.get('width', 1.0) * img_width
                                    bbox_height = image_bbox.get('height', 1.0) * img_height

                                    # 计算裁剪区域
                                    left = max(0, int(center_x - bbox_width / 2))
                                    top = max(0, int(center_y - bbox_height / 2))
                                    right = min(img_width, int(center_x + bbox_width / 2))
                                    bottom = min(img_height, int(center_y + bbox_height / 2))

                                    # 裁剪图片
                                    cropped_img = img.crop((left, top, right, bottom))

                                    # 保存裁剪后的图片到临时文件
                                    import tempfile
                                    temp_dir = tempfile.mkdtemp(prefix=f"cropped_images_{uuid.uuid4().hex[:8]}_")
                                    cropped_filename = f"cropped_{os.path.splitext(filename)[0]}.{image_ext if 'image_ext' in locals() else 'jpg'}"
                                    cropped_path = os.path.join(temp_dir, cropped_filename)

                                    # 保存为JPEG格式
                                    if cropped_img.mode in ('RGBA', 'LA', 'P'):
                                        # 转换为RGB
                                        background = PILImage.new('RGB', cropped_img.size, (255, 255, 255))
                                        if cropped_img.mode == 'P':
                                            cropped_img = cropped_img.convert('RGBA')
                                        background.paste(cropped_img, mask=cropped_img.split()[-1] if cropped_img.mode == 'RGBA' else None)
                                        cropped_img = background
                                    elif cropped_img.mode != 'RGB':
                                        cropped_img = cropped_img.convert('RGB')

                                    cropped_img.save(cropped_path, 'JPEG', quality=95)

                                    # 🚨 为裁剪图片生成独立的UUID
                                    cropped_uuid = str(uuid.uuid4())

                                    # 添加裁剪图片信息到结果
                                    result['cropped_image_uuid'] = cropped_uuid
                                    result['cropped_image_path'] = cropped_path
                                    result['cropped_image_available'] = True
                                    result['cropped_temp_dir'] = temp_dir
                                    result['cropped_image_filename'] = cropped_filename

                                    logger.info(f"成功裁剪医学影像: {cropped_path}, UUID: {cropped_uuid}, 尺寸: {right-left}x{bottom-top}")

                                except Exception as crop_error:
                                    logger.error(f"裁剪医学影像失败: {str(crop_error)}")
                                    result['cropped_image_available'] = False
                        else:
                            logger.warning(f"图片内容提取为空: {filename}")
                            result['file_content'] = f"图片文件: {filename} (API调用成功但返回内容为空)"
                            result['has_medical_image'] = False
                            result['model_used'] = successful_model
                            result['extraction_error'] = 'API调用成功但返回内容为空'
                    except Exception as json_error:
                        logger.warning(f"JSON解析失败: {filename}, 错误: {str(json_error)}, 原始内容: {response_content[:200]}")
                        # 降级处理：直接使用原始内容
                        image_description = response_content.strip()
                        if image_description.startswith("\n<|begin_of_box|>"):
                            image_description = image_description[len("\n<|begin_of_box|>"):].strip()
                        elif image_description.startswith("<|begin_of_box|>"):
                            image_description = image_description[len("<|begin_of_box|>"):].strip()
                        if image_description.endswith("<|end_of_box|>"):
                            image_description = image_description[:-len("<|end_of_box|>")].strip()

                        result['file_content'] = image_description
                        result['has_medical_image'] = False
                        result['model_used'] = successful_model
                        result['extraction_success'] = True if image_description else False  # 只有非空内容才算成功
                        if not result['extraction_success']:
                            result['extraction_error'] = 'JSON解析失败且内容为空'
                else:
                    logger.warning(f"图片内容提取为空: {filename}")
                    result['file_content'] = f"图片文件: {filename} (API调用成功但返回内容为空)"
                    result['has_medical_image'] = False
                    result['model_used'] = successful_model

                return result

            except Exception as outer_error:
                # 构建详细的错误信息
                error_details = str(outer_error)
                error_type = type(outer_error).__name__
                status_code = None
                response_text = None

                # 尝试提取HTTP响应详情
                if hasattr(outer_error, 'response') and outer_error.response is not None:
                    try:
                        status_code = getattr(outer_error.response, 'status_code', None)
                        if hasattr(outer_error.response, 'json'):
                            error_json = outer_error.response.json()
                            response_text = f"API错误响应: {error_json}"
                        elif hasattr(outer_error.response, 'text'):
                            response_text = outer_error.response.text[:500]  # 限制长度
                    except Exception as parse_error:
                        response_text = f"无法解析响应: {str(parse_error)}"

                # 记录详细的错误日志
                logger.error(f"""多模态API调用失败详情:
文件名: {filename}
错误类型: {error_type}
错误信息: {error_details}
HTTP状态码: {status_code}
响应内容: {response_text}
文件格式: {actual_format or file_extension}
MIME类型: {mime_type}
尝试的模型: {multimodal_model_names}
文件大小: {file_size / 1024:.1f} KB""")

                # 构建返回的错误信息
                result['file_content'] = f"""图片文件: {filename}
处理状态: 多模态API调用失败
错误类型: {error_type}
错误信息: {error_details}
HTTP状态码: {status_code if status_code else 'N/A'}
文件格式: {actual_format or file_extension}
MIME类型: {mime_type}
尝试的模型: {', '.join(multimodal_model_names)}"""
                result['extraction_error'] = f"{error_type}: {error_details[:200]}"  # 截取前200字符作为错误信息

                return result

            finally:
                # 清理临时转换文件
                if converted_file and os.path.exists(converted_file):
                    try:
                        # 删除转换后的文件
                        os.unlink(converted_file)
                        # 删除临时目录
                        temp_dir = os.path.dirname(converted_file)
                        if os.path.exists(temp_dir):
                            os.rmdir(temp_dir)
                        logger.debug(f"清理HEIC转换临时文件: {converted_file}")
                    except Exception as cleanup_error:
                        logger.warning(f"清理临时文件失败: {cleanup_error}")

        except Exception as e:
            logger.error(f"""图片处理整体失败:
文件名: {filename if 'filename' in locals() else (os.path.basename(original_path) if 'original_path' in locals() else os.path.basename(path))}
文件路径: {original_path if 'original_path' in locals() else path}
错误类型: {type(e).__name__}
错误信息: {str(e)}
堆栈跟踪: {str(e.__traceback__) if hasattr(e, '__traceback__') else 'N/A'}""")
            filename = os.path.basename(original_path) if 'original_path' in locals() else (os.path.basename(path) if 'filename' not in locals() else filename)
            result = {
                'file_extension': 'image',
                'file_name': filename,
                'extraction_success': False,
                'extraction_error': f"{type(e).__name__}: {str(e)}"
            }
            result['file_content'] = f"图片文件: {filename} (处理失败: {type(e).__name__}: {str(e)})"

            # 确保清理临时文件
            if 'converted_file' in locals() and converted_file and os.path.exists(converted_file):
                try:
                    os.unlink(converted_file)
                    temp_dir = os.path.dirname(converted_file)
                    if os.path.exists(temp_dir):
                        os.rmdir(temp_dir)
                except:
                    pass

            return result

    def read_md(self, path):
        try:
            filename = os.path.basename(path)
            result = {'file_extension': 'md', 'file_name': filename}
            with open(path, 'r', encoding='utf-8') as file:
                result['file_content'] = file.read()
            return result
        except Exception as e:
            logger.error(f"""MD文件读取失败:
文件名: {filename if 'filename' in locals() else os.path.basename(path)}
文件路径: {path}
错误类型: {type(e).__name__}
错误信息: {str(e)}""")
            return {'file_extension': 'md', 'file_name': filename, 'file_content': ''}

    def read_json(self, path):
        try:
            filename = os.path.basename(path)
            result = {'file_extension': 'json', 'file_name': filename}
            with open(path, 'r', encoding='utf-8') as file:
                content = json.load(file)
                # Convert JSON object to string to maintain consistency
                result['file_content'] = content
            return result
        except Exception as e:
            logger.error(f"""JSON文件读取失败:
文件名: {filename if 'filename' in locals() else os.path.basename(path)}
文件路径: {path}
错误类型: {type(e).__name__}
错误信息: {str(e)}""")
            return {'file_extension': 'json', 'file_name': filename, 'file_content': ''}

    def extract_text_html2text(self, url, ignore_links=False, ignore_images=False, ignore_videos=False):
        try:
            response = requests.get(url)
            response.raise_for_status()
            html_content = response.text

            h = html2text.HTML2Text()
            h.ignore_links = ignore_links
            h.ignore_images = ignore_images
            h.ignore_videos = ignore_videos
            text_content = h.handle(html_content)

            return text_content
        except Exception as e:
            logger.info(f"html2text extraction failed: {e}")
            return ""

    def extract_content_goose(self, url):
        try:
            g = Goose({'stopwords_class': StopWordsChinese})
            article = g.extract(url=url)
            html = article.raw_html
            soup = BeautifulSoup(html, 'html.parser')

            for img in soup.find_all('img'):
                img_src = img.get('src')
                if img_src:
                    img.replace_with(f'[Image: {img_src}]')

            for video in soup.find_all('video'):
                video_src = video.get('src')
                if video_src:
                    video.replace_with(f'[Video: {video_src}]')

            for a in soup.find_all('a'):
                link_href = a.get('href')
                if link_href:
                    a.replace_with(f'[Link: {link_href}]')

            for code in soup.find_all('pre'):
                code_text = code.get_text()
                if code_text:
                    code.replace_with(f'[Code Block: {code_text}]')

            text = soup.get_text()
            lines = text.splitlines()
            cleaned_lines = [line.strip() for line in lines if line.strip()]
            cleaned_text = "\n".join(cleaned_lines)

            return cleaned_text
        except Exception as e:
            logger.info(f"Goose extraction failed: {e}")
            return ""

    def extract_content_trafilatura(self, url):
        try:
            downloaded = trafilatura.fetch_url(url)
            content = trafilatura.extract(downloaded, include_images=True, include_links=True, include_formatting=True)
            return content
        except Exception as e:
            logger.info(f"Trafilatura extraction failed: {e}")
            return ""

    def read_url(self, url, priority=None):
        methods = {
            'Trafilatura': self.extract_content_trafilatura,
            'html2text': self.extract_text_html2text,
            'Goose': self.extract_content_goose
        }
        
        if not priority:
            priority = list(methods.keys())
            random.shuffle(priority)
        
        for method_name in priority:
            method = methods[method_name]
            text = method(url)
            if text and len(text) > 50 and "当前环境异常，完成验证后即可继续访问" not in text:
                return {'file_extension': 'url', 'file_name': url, 'file_content': text}

        return None

    def unzip_file(self, zip_file_path):
        extract_to_path = os.path.dirname(zip_file_path)
        zip_filename_without_ext = os.path.splitext(os.path.basename(zip_file_path))[0]
        extract_path = os.path.join(extract_to_path, zip_filename_without_ext)
        
        os.makedirs(extract_path, exist_ok=True)
        
        try:
            with zipfile.ZipFile(zip_file_path, 'r') as zip_ref:
                zip_info_list = zip_ref.infolist()
                total_files = len(zip_info_list)
                logger.info(f"开始解压ZIP文件，包含 {total_files} 个条目")
                
                extracted_count = 0
                skipped_count = 0
                
                for zip_info in zip_info_list:
                    try:
                        original_filename = zip_info.filename
                        
                        # 尝试不同的编码方式解码文件名
                        filename = None
                        encoding_attempts = [
                            ('cp437', 'utf-8'),
                            ('cp437', 'gbk'),
                            ('utf-8', None),
                        ]
                        
                        for src_encoding, target_encoding in encoding_attempts:
                            try:
                                if target_encoding:
                                    filename = original_filename.encode(src_encoding).decode(target_encoding)
                                else:
                                    detected_encoding = chardet.detect(original_filename.encode('utf-8'))
                                    detected_enc = detected_encoding.get('encoding', 'utf-8')
                                    filename = original_filename.encode('utf-8').decode(detected_enc)
                                break
                            except (UnicodeDecodeError, TypeError):
                                continue
                        
                        if filename is None:
                            logger.warning(f"文件名编码失败，跳过: {original_filename}")
                            skipped_count += 1
                            continue

                        zip_info.filename = filename
                        zip_ref.extract(zip_info, extract_path)
                        extracted_count += 1
                        
                    except Exception as e:
                        logger.warning(f"解压文件失败: {zip_info.filename}")
                        skipped_count += 1
                        continue
                
                logger.info(f"ZIP解压完成: 成功 {extracted_count} 个, 跳过 {skipped_count} 个")
                
        except Exception as e:
            logger.error(f"解压ZIP文件失败: {str(e)}")
            raise
        
        return extract_path

    def process_zip_file(self, zip_file_path, processed_zips=None, cleanup=True):
        """
        处理zip文件：解压并递归提取所有支持的文件类型内容，包括嵌套的zip文件
        返回格式与其他文件类型保持一致
        """
        start_time = time.time()
        zip_filename = os.path.basename(zip_file_path)
        logger.info(f"开始处理ZIP文件: {zip_filename}")
        
        # 防止循环引用的zip文件
        if processed_zips is None:
            processed_zips = set()
        
        # 获取zip文件的绝对路径，避免重复处理
        abs_zip_path = os.path.abspath(zip_file_path)
        if abs_zip_path in processed_zips:
            logger.warning(f"检测到循环引用的zip文件，跳过: {zip_filename}")
            return []
        
        processed_zips.add(abs_zip_path)
        
        extract_path = None
        persistent_temp_dir = None  # 🚨 新增：持久临时目录，用于保存需要上传的文件
        
        try:
            file_size_mb = os.path.getsize(zip_file_path) / 1024 / 1024
            logger.info(f"ZIP文件大小: {file_size_mb:.1f} MB")
            
            # 创建持久临时目录，用于保存需要上传的二进制文件
            import tempfile
            persistent_temp_dir = tempfile.mkdtemp(prefix=f"zip_upload_{uuid.uuid4().hex[:8]}_")
            
            # 解压zip文件
            extract_path = self.unzip_file(zip_file_path)
            
            if not os.path.exists(extract_path):
                logger.error(f"解压失败，路径不存在: {extract_path}")
                return [{
                    'file_extension': 'zip',
                    'file_name': zip_filename,
                    'file_content': f'解压失败：路径不存在'
                }]
            
            # 获取解压后的所有文件
            extracted_files = self.get_all_files_recursive(extract_path)
            logger.info(f"ZIP解压完成，包含 {len(extracted_files)} 个文件")
            
            if not extracted_files:
                logger.warning(f"ZIP文件解压后未找到任何文件")
                return [{
                    'file_extension': 'zip',
                    'file_name': zip_filename,
                    'file_content': '解压后未找到任何文件'
                }]
            
            # 分类处理文件
            non_zip_files = []
            zip_files = []
            other_files = []
            
            for file_path in extracted_files:
                file_ext = os.path.splitext(file_path)[1].lower()
                
                if file_ext == '.zip':
                    zip_files.append(file_path)
                elif file_ext in ['.pdf', '.docx', '.pptx', '.ppt', '.txt', '.md', '.json', '.png', '.jpg', '.jpeg', '.webp', '.heic', '.heif']:
                    non_zip_files.append(file_path)
                else:
                    other_files.append(file_path)
            
            logger.info(f"文件分类: {len(non_zip_files)} 个支持文件, {len(zip_files)} 个嵌套ZIP, {len(other_files)} 个其他文件")
            
            # 处理所有文件
            processed_files = []
            
            # 处理非zip文件
            if non_zip_files:
                # 从环境变量获取并发数配置
                max_workers = int(os.getenv("PDF_IMAGE_CONCURRENT_WORKERS", "10"))
                logger.info(f"开始并发处理 {len(non_zip_files)} 个文件（并发数：{max_workers}）...")
                success_count = 0
                fail_count = 0

                # 使用线程池并发处理
                with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
                    futures = {executor.submit(self._run_single_file, file_path): file_path for file_path in non_zip_files}
                    
                    for future in concurrent.futures.as_completed(futures):
                        file_path = futures[future]
                        filename = os.path.basename(file_path)
                        
                        try:
                            result = future.result()
                            if result:
                                if isinstance(result, list):
                                    for sub_result in result:
                                        if isinstance(sub_result, dict):
                                            original_path = file_path
                                            file_ext = os.path.splitext(filename)[1].lower()
                                            
                                            if file_ext in ['.png', '.jpg', '.jpeg', '.webp', '.heic', '.heif', '.pdf', '.docx', '.pptx', '.ppt']:
                                                # 为二进制文件创建持久副本
                                                persistent_filename = f"{sub_result.get('file_uuid', str(uuid.uuid4()))}{file_ext}"
                                                persistent_path = os.path.join(persistent_temp_dir, persistent_filename)
                                                
                                                try:
                                                    shutil.copy2(original_path, persistent_path)
                                                    sub_result['original_file_path'] = persistent_path
                                                    sub_result['temp_file_available'] = True
                                                    sub_result['persistent_temp_file'] = True
                                                except Exception as e:
                                                    logger.warning(f"复制文件失败: {filename}")
                                                    sub_result['original_file_path'] = original_path
                                                    sub_result['temp_file_available'] = False
                                            else:
                                                sub_result['original_file_path'] = original_path
                                                sub_result['temp_file_available'] = True
                                    
                                    processed_files.extend(result)
                                    success_count += len(result)
                                else:
                                    if isinstance(result, dict):
                                        original_path = file_path
                                        file_ext = os.path.splitext(filename)[1].lower()
                                        
                                        if file_ext in ['.png', '.jpg', '.jpeg', '.webp', '.heic', '.heif', '.pdf', '.docx', '.pptx', '.ppt']:
                                            persistent_filename = f"{result.get('file_uuid', str(uuid.uuid4()))}{file_ext}"
                                            persistent_path = os.path.join(persistent_temp_dir, persistent_filename)
                                            
                                            try:
                                                shutil.copy2(original_path, persistent_path)
                                                result['original_file_path'] = persistent_path
                                                result['temp_file_available'] = True
                                                result['persistent_temp_file'] = True
                                            except Exception as e:
                                                logger.warning(f"复制文件失败: {filename}")
                                                result['original_file_path'] = original_path
                                                result['temp_file_available'] = False
                                        else:
                                            result['original_file_path'] = original_path
                                            result['temp_file_available'] = True
                                    
                                    processed_files.append(result)
                                    success_count += 1
                            else:
                                fail_count += 1
                                logger.warning(f"文件处理返回空结果: {filename}")
                        except Exception as e:
                            fail_count += 1
                            logger.error(f"处理文件失败 {filename}: {str(e)}")
                            continue
                
                logger.info(f"文件处理完成: 成功 {success_count} 个, 失败 {fail_count} 个")
            
            # 处理嵌套zip文件
            if zip_files:
                logger.info(f"开始处理 {len(zip_files)} 个嵌套ZIP文件...")
                nested_success_count = 0
                nested_fail_count = 0
                
                # 使用较小的线程池处理嵌套zip文件
                with concurrent.futures.ThreadPoolExecutor(max_workers=2) as executor:
                    process_func = partial(self.process_zip_file, processed_zips=processed_zips.copy(), cleanup=False)
                    futures = {executor.submit(process_func, zip_file_path): zip_file_path for zip_file_path in zip_files}
                    
                    for future in concurrent.futures.as_completed(futures):
                        zip_file_path = futures[future]
                        zip_filename = os.path.basename(zip_file_path)
                        
                        try:
                            nested_results = future.result()
                            
                            if nested_results:
                                if isinstance(nested_results, list):
                                    processed_files.extend(nested_results)
                                    nested_success_count += len(nested_results)
                                else:
                                    processed_files.append(nested_results)
                                    nested_success_count += 1
                                logger.info(f"嵌套ZIP处理完成: {zip_filename}")
                            else:
                                nested_fail_count += 1
                                logger.warning(f"嵌套ZIP处理返回空结果: {zip_filename}")
                        except Exception as e:
                            nested_fail_count += 1
                            logger.error(f"处理嵌套ZIP失败 {zip_filename}: {str(e)}")
                            continue
                
                logger.info(f"嵌套ZIP处理完成: 成功 {nested_success_count} 个文件, 失败 {nested_fail_count} 个ZIP")
            
            total_elapsed = time.time() - start_time
            logger.info(f"ZIP文件处理完成: {zip_filename} (用时 {total_elapsed:.1f}s, 共 {len(processed_files)} 个文件)")
            
            # 为结果添加清理信息
            if persistent_temp_dir and processed_files:
                for pf in processed_files:
                    if isinstance(pf, dict) and pf.get('persistent_temp_file'):
                        pf['cleanup_temp_dir'] = persistent_temp_dir
            
            return processed_files
            
        except Exception as e:
            total_elapsed = time.time() - start_time
            logger.error(f"ZIP文件处理失败: {zip_filename} - {str(e)} (用时 {total_elapsed:.1f}s)")
            return [{
                'file_extension': 'zip',
                'file_name': zip_filename,
                'file_content': f'ZIP文件处理失败: {str(e)}'
            }]
        finally:
            # 清理解压的临时文件夹（只有顶级zip才清理）
            if cleanup and extract_path and os.path.exists(extract_path):
                try:
                    shutil.rmtree(extract_path)
                    logger.debug(f"清理临时文件夹: {extract_path}")
                except Exception as e:
                    logger.warning(f"清理临时文件夹失败: {str(e)}")
            
            # 注意：不清理persistent_temp_dir，让conversations.py在上传完成后清理

    def get_all_files_recursive(self, folder_path):
        """递归获取文件夹中的所有文件路径"""
        file_paths = []
        
        try:
            for root, dirs, files in os.walk(folder_path):
                for file in files:
                    file_path = os.path.join(root, file)
                    file_paths.append(file_path)
            
            logger.debug(f"递归扫描完成: {len(file_paths)} 个文件")
            
        except Exception as e:
            logger.error(f"递归遍历文件夹失败: {str(e)}")
        
        return file_paths

    def _run_single_file(self, path: str):
        """处理单个文件，不递归处理zip文件"""
        filename = os.path.basename(path)
        
        # 过滤系统隐藏文件
        if filename.startswith('._') or filename.startswith('.DS_Store'):
            logger.warning(f"跳过系统隐藏文件: {filename}")
            return {
                'file_extension': 'hidden',
                'file_name': filename,
                'file_content': f"系统隐藏文件: {filename} (已跳过处理)"
            }
        
        if not os.path.isfile(path):
            logger.warning(f"路径不是有效文件: {path}")
            return None
            
        file_extension = os.path.splitext(path)[1].lower()
        
        try:
            # 处理各种文件类型，但不处理zip文件
            result = None
            if file_extension == '.pdf':
                result = self.read_pdf(path)
            elif file_extension == '.docx':
                result = self.read_docx(path)
            elif file_extension in ['.pptx', '.ppt']:
                result = self.read_ppt(path)
            elif file_extension == '.txt':
                result = self.read_txt(path)
            elif file_extension == '.md':
                result = self.read_md(path)
            elif file_extension == '.json':
                result = self.read_json(path)
            elif file_extension in ['.png', '.jpg', '.jpeg', '.webp', '.heic', '.heif']:
                result = self.read_image(path)
            elif file_extension == '.zip':
                # 在单文件处理中不处理zip，避免重复
                return None
            else:
                result = self.read_file(path)
            
            if result and isinstance(result, dict):
                content = result.get('file_content', '')
                content_length = len(str(content)) if content else 0
                if content_length > 0:
                    logger.debug(f"文件处理成功: {filename} ({content_length} 字符)")
                else:
                    logger.warning(f"文件处理完成但内容为空: {filename}")
            
            return result
            
        except Exception as e:
            logger.error(f"""单文件处理失败:
文件名: {filename}
文件路径: {path}
文件类型: {file_extension}
错误类型: {type(e).__name__}
错误信息: {str(e)}""")
            return None

    def list_files_in_folder(self, folder_path):
        file_paths = []
        for filename in os.listdir(folder_path):
            filepath = os.path.join(folder_path, filename)
            if os.path.isfile(filepath):
                file_paths.append(filepath)
        return file_paths

    def read_file(self, file_path):
        result = {}
        file_name = os.path.basename(file_path)
        file_name_without_extension, file_extension = os.path.splitext(file_name)

        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                file_content = file.read()
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='gbk') as file:
                    file_content = file.read()
            except Exception:
                return []
        except Exception:
            return []

        result['file_extension'] = file_extension
        result['file_name'] = file_name
        result['file_content'] = file_content

        return result


    def is_text_garbled(self, text):
        chinese_characters = re.findall(r'[\u4e00-\u9fff]', text)
        symbol_characters = re.findall(r'[\u0000-\u0020\u3000\uFFFD]', text)

        if len(chinese_characters) > 0:
            chinese_ratio = len(chinese_characters) / max(len(text), 1)
            symbol_ratio = len(symbol_characters) / max(len(text), 1)
            return chinese_ratio < 0.2 or symbol_ratio > 0.3

        non_ascii_ratio = sum(1 for char in text if ord(char) > 127) / max(len(text), 1)
        return non_ascii_ratio > 0.3

    def read_pdf_with_images(self, filepath):
        """
        新的PDF处理方式：提取文本和图片，图片用多模态模型处理
        返回列表：[PDF本身, 图片1, 图片2, ...]
        """
        filename = os.path.basename(filepath)
        pdf_uuid = str(uuid.uuid4())

        logger.info(f"开始提取PDF内容（带图片模式）: {filename}")

        def clean_text(text):
            try:
                return text.encode('utf-8', 'ignore').decode('utf-8')
            except UnicodeDecodeError:
                try:
                    return text.encode('utf-8', 'ignore').decode('gbk')
                except UnicodeDecodeError:
                    return text

        try:
            import tempfile

            # 打开PDF文档
            doc = fitz.open(filepath)

            full_content_parts = []  # 完整内容（文本+图片描述按页面顺序）
            extracted_images = []  # 提取的图片列表
            total_images = 0

            # 创建临时目录保存提取的图片
            temp_dir = tempfile.mkdtemp(prefix=f"pdf_images_{uuid.uuid4().hex[:8]}_")
            logger.info(f"创建临时目录: {temp_dir}")

            try:
                # 保存页面总数（在关闭文档前）
                total_pages = len(doc)

                # 步骤1：先渲染所有页面为图片（不调用多模态模型）
                page_image_paths = []
                for page_num in range(total_pages):
                    page = doc[page_num]
                    page_number = page_num + 1

                    # 1. 提取页面文本
                    page_text = clean_text(page.get_text())
                    if page_text and page_text.strip():
                        full_content_parts.append(f"[第{page_number}页]\n{page_text.strip()}")
                        logger.info(f"提取第{page_number}页文本: {len(page_text)} 字符")

                    # 2. 渲染页面为图片
                    try:
                        pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
                        img_filename = f"page_{page_number}.png"
                        img_path = os.path.join(temp_dir, img_filename)
                        pix.save(img_path)

                        width = pix.width
                        height = pix.height
                        file_size = os.path.getsize(img_path)

                        logger.info(f"渲染第{page_number}页为完整图片: {img_filename} ({file_size} bytes, {width}x{height}px)")
                        page_image_paths.append((page_number, img_path))

                    except Exception as img_error:
                        logger.error(f"渲染第{page_number}页为图片时出错: {str(img_error)}")
                        continue

                # 关闭PDF文档
                doc.close()

                # 步骤2：并发调用多模态模型处理所有图片
                if page_image_paths:
                    max_workers = int(os.getenv("PDF_IMAGE_CONCURRENT_WORKERS", "10"))
                    logger.info(f"开始并发处理 {len(page_image_paths)} 张PDF页面图片（并发数：{max_workers}）...")

                    def process_page_image(page_info):
                        """处理单个页面图片"""
                        page_number, img_path = page_info
                        try:
                            image_result = self.read_image(img_path)
                            return (page_number, img_path, image_result)
                        except Exception as e:
                            logger.error(f"处理第{page_number}页图片失败: {str(e)}")
                            return (page_number, img_path, None)

                    # 使用线程池并发处理
                    with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
                        futures = {executor.submit(process_page_image, page_info): page_info for page_info in page_image_paths}

                        for future in concurrent.futures.as_completed(futures):
                            page_info = futures[future]
                            page_number = page_info[0]

                            try:
                                page_number, img_path, image_result = future.result()

                                if image_result:
                                    img_uuid = str(uuid.uuid4())
                                    image_description = image_result.get('file_content', '')
                                    has_medical_image = image_result.get('has_medical_image', False)

                                    # 插入图片描述到完整内容中（保持顺序）
                                    if image_description and image_description.strip():
                                        full_content_parts.append(
                                            f"\n[第{page_number}页 - 完整页面图片]\n{image_description.strip()}\n"
                                        )
                                        logger.info(f"提取第{page_number}页完整图片内容: {len(image_description)} 字符, 医学影像: {has_medical_image}")

                                    # 构建图片项
                                    image_item = {
                                        'file_uuid': img_uuid,
                                        'file_name': f"{os.path.splitext(filename)[0]}_第{page_number}页.png",
                                        'file_extension': 'png',
                                        'file_content': image_description,
                                        'has_medical_image': has_medical_image,
                                        'model_used': image_result.get('model_used'),

                                        # 来源信息
                                        'source_type': 'rendered_pdf_page',
                                        'parent_pdf_uuid': pdf_uuid,
                                        'parent_pdf_filename': filename,

                                        # 位置信息
                                        'page_number': page_number,
                                        'image_index_in_page': 0,
                                        'position_in_parent': f"page_{page_number}_full",

                                        # 临时文件路径
                                        'temp_file_path': img_path,
                                        'temp_file_available': True,
                                        'persistent_temp_file': True,
                                        'cleanup_temp_dir': temp_dir,

                                        # 医学影像裁剪信息
                                        'image_bbox': image_result.get('image_bbox'),
                                        'cropped_image_uuid': image_result.get('cropped_image_uuid'),
                                        'cropped_image_path': image_result.get('cropped_image_path'),
                                        'cropped_image_filename': image_result.get('cropped_image_filename'),
                                        'cropped_image_available': image_result.get('cropped_image_available', False),
                                        'cropped_temp_dir': image_result.get('cropped_temp_dir')
                                    }

                                    extracted_images.append(image_item)
                                    total_images += 1
                                else:
                                    logger.warning(f"第{page_number}页完整图片提取失败")

                            except Exception as e:
                                logger.error(f"处理第{page_number}页结果时出错: {str(e)}")
                                continue

                # 构建PDF本身的项
                pdf_item = {
                    'file_uuid': pdf_uuid,
                    'file_name': filename,
                    'file_extension': 'pdf',
                    'file_content': '\n\n'.join(full_content_parts) if full_content_parts else '',
                    'extraction_mode': 'with_images',
                    'extracted_image_count': total_images,
                    'extraction_success': True
                }

                logger.info(f"PDF提取完成: {filename}, 总页数: {total_pages}, 提取图片: {total_images} 张")

                # 返回列表：[PDF本身, 图片1, 图片2, ...]
                return [pdf_item] + extracted_images

            except Exception as e:
                logger.error(f"PDF图片提取过程出错: {filename}, 错误: {str(e)}")
                # 清理临时目录
                if os.path.exists(temp_dir):
                    try:
                        shutil.rmtree(temp_dir)
                    except:
                        pass
                raise

        except Exception as e:
            logger.error(f"PDF图片提取整体失败: {filename}, 错误: {str(e)}")
            return {
                'file_uuid': pdf_uuid,
                'file_name': filename,
                'file_extension': 'pdf',
                'file_content': f'PDF图片提取失败: {str(e)}',
                'extraction_success': False,
                'extraction_error': f"{type(e).__name__}: {str(e)}"
            }

    def read_pdf(self, filepath):
        """原有的PDF处理方式：只提取文本"""
        filename = os.path.basename(filepath)
        result = {'file_extension': 'pdf', 'file_name': filename}

        def clean_text(text):
            # 尝试先用 UTF-8 解码，如果失败则使用 GBK 解码
            try:
                return text.encode('utf-8', 'ignore').decode('utf-8')
            except UnicodeDecodeError:
                try:
                    return text.encode('utf-8', 'ignore').decode('gbk')
                except UnicodeDecodeError:
                    return text  # 如果两种解码都失败，保持原始文本

        # 1. 尝试使用 OCR API
        # 定义新的OCR API URL
        ocr_api_url = 'https://pdf-parsing.dev.6ccloud.com/pdf_parsing_api'
        try:
            with open(filepath, 'rb') as file:
                pdf = pypdf.PdfReader(file)
                num_pages = len(pdf.pages)
                all_markdown_content = []
                
                for i in range(0, num_pages, 20):
                    end_page = min(i + 20, num_pages)
                    
                    # 创建一个新的 PDF 写入器
                    pdf_writer = pypdf.PdfWriter()
                    for page_num in range(i, end_page):
                        pdf_writer.add_page(pdf.pages[page_num])
                    
                    # 将切分后的 PDF 保存到内存中
                    temp_pdf = io.BytesIO()
                    pdf_writer.write(temp_pdf)
                    temp_pdf.seek(0)
                    
                    # 发送请求到新的 OCR API
                    files = {'pdf_file': (f'{filename}_part_{i//20+1}.pdf', temp_pdf, 'application/pdf')}
                    response = requests.post(ocr_api_url, files=files, headers={'accept': 'application/json'}, timeout=600)
                    
                    # 处理响应
                    if response.status_code == 200:
                        response_json = response.json()
                        markdown_content = response_json.get('markdown', '')
                        if markdown_content:
                            all_markdown_content.append(markdown_content)
                        logger.info(f"{filename} part {i//20+1}, OCR API returned markdown content: {markdown_content[:50]}...")
                    elif response.status_code in [400, 404, 420, 500]:
                        # 记录详细的错误信息
                        try:
                            error_body = response.json() if response.content else "无响应内容"
                        except:
                            error_body = response.text[:500] if response.text else "无法解析响应"
                        logger.error(f"""OCR API调用失败详情:
文件名: {filename} part {i//20+1}
HTTP状态码: {response.status_code}
响应内容: {error_body}
请求URL: {ocr_api_url}""")
                        # 继续尝试其他页
                    else:
                        # 记录详细的非预期状态码错误
                        try:
                            error_body = response.json() if response.content else "无响应内容"
                        except:
                            error_body = response.text[:500] if response.text else "无法解析响应"
                        logger.error(f"""OCR API返回非预期状态码:
文件名: {filename} part {i//20+1}
HTTP状态码: {response.status_code}
响应内容: {error_body}
请求URL: {ocr_api_url}""")
                
                if all_markdown_content:
                    combined_ocr_text = clean_text("".join(all_markdown_content))
                    logger.info(f"OCR API extracted: {combined_ocr_text[:50]}...")  # 打印前 50 个字符
                    if combined_ocr_text and not self.is_text_garbled(combined_ocr_text):
                        result['file_content'] = combined_ocr_text
                        return result
                logger.info(f"OCR result is empty or garbled for {filename}")
            
        except Exception as e:
            logger.error(f"""OCR API整体处理失败:
文件名: {filename}
错误类型: {type(e).__name__}
错误信息: {str(e)}
请求URL: {ocr_api_url}""")
            if 'response' in locals():
                try:
                    logger.error(f"""OCR响应详情:
状态码: {response.status_code}
响应内容: {response.text[:500]}""")
                except:
                    pass

        # 2. 尝试使用 MinerU
        mineru_token = os.getenv("MINERU_TOKEN")
        if mineru_token:
            try:
                logger.info(f"Attempting to extract {filename} using MinerU...")

                # MinerU API配置
                mineru_base_url = "https://mineru.net/api/v4"

                # Step 1: 申请上传URL
                upload_url_endpoint = f"{mineru_base_url}/file-urls/batch"
                headers = {
                    "Content-Type": "application/json",
                    "Authorization": f"Bearer {mineru_token}"
                }
                upload_data = {
                    "enable_formula": True,
                    "language": "ch",
                    "enable_table": True,
                    "files": [
                        {
                            "name": filename,
                            "is_ocr": True,
                            "data_id": f"pdf_{uuid.uuid4().hex[:8]}"
                        }
                    ]
                }

                upload_response = requests.post(upload_url_endpoint, headers=headers, json=upload_data, timeout=30)

                if upload_response.status_code == 200:
                    upload_result = upload_response.json()
                    if upload_result.get("code") == 0:
                        batch_id = upload_result["data"]["batch_id"]
                        upload_url = upload_result["data"]["file_urls"][0]
                        logger.info(f"MinerU: Upload URL obtained, batch_id={batch_id}")

                        # Step 2: 上传文件
                        with open(filepath, 'rb') as f:
                            upload_file_response = requests.put(upload_url, data=f, timeout=300)

                        if upload_file_response.status_code == 200:
                            logger.info(f"MinerU: File uploaded successfully")

                            # Step 3: 等待处理并轮询结果
                            time.sleep(5)  # 等待系统处理
                            result_url = f"{mineru_base_url}/extract-results/batch/{batch_id}"
                            max_retries = 60
                            retry_interval = 5

                            for attempt in range(max_retries):
                                result_response = requests.get(result_url, headers=headers, timeout=30)

                                if result_response.status_code == 200:
                                    result_data = result_response.json()

                                    if result_data.get("code") == 0:
                                        extract_results = result_data["data"]["extract_result"]

                                        for file_result in extract_results:
                                            state = file_result["state"]

                                            if state == "done":
                                                # Step 4: 下载结果ZIP
                                                zip_url = file_result.get("full_zip_url")
                                                if zip_url:
                                                    logger.info(f"MinerU: Extraction completed, downloading result...")

                                                    # 创建临时目录保存结果
                                                    output_dir = Path("/home/ubuntu/github/mediwise/output/mineru")
                                                    output_dir.mkdir(parents=True, exist_ok=True)

                                                    zip_filename = f"{uuid.uuid4().hex[:8]}_{filename.replace('.pdf', '')}_result.zip"
                                                    zip_path = output_dir / zip_filename

                                                    # 下载ZIP文件
                                                    zip_response = requests.get(zip_url, stream=True, timeout=300)
                                                    if zip_response.status_code == 200:
                                                        with open(zip_path, 'wb') as f:
                                                            for chunk in zip_response.iter_content(chunk_size=8192):
                                                                f.write(chunk)

                                                        logger.info(f"MinerU: Result downloaded to {zip_path}")

                                                        # Step 5: 解压并提取MD内容
                                                        extract_folder = output_dir / f"{uuid.uuid4().hex[:8]}_{filename.replace('.pdf', '')}_extracted"
                                                        extract_folder.mkdir(exist_ok=True)

                                                        try:
                                                            with zipfile.ZipFile(zip_path, 'r') as zip_ref:
                                                                zip_ref.extractall(extract_folder)

                                                            # 查找MD文件
                                                            md_files = list(extract_folder.rglob("*.md"))
                                                            if md_files:
                                                                # 读取第一个MD文件的内容
                                                                with open(md_files[0], 'r', encoding='utf-8') as md_file:
                                                                    md_content = md_file.read()

                                                                if md_content:
                                                                    md_content = clean_text(md_content)
                                                                    logger.info(f"MinerU extracted: {md_content[:50]}...")
                                                                    result['file_content'] = md_content
                                                                    logger.info(f"MinerU extraction successful for {filename}")
                                                                    return result
                                                                else:
                                                                    logger.info(f"MinerU result is empty for {filename}")
                                                                    break  # 空内容，跳出重试循环
                                                            else:
                                                                logger.warning(f"MinerU: No MD files found in extracted result")
                                                                break  # 没有找到MD文件，跳出重试循环

                                                        except Exception as extract_error:
                                                            logger.error(f"MinerU: Error extracting ZIP: {extract_error}")
                                                            break  # 解压出错，跳出重试循环
                                                    else:
                                                        logger.error(f"MinerU: Failed to download result ZIP, status={zip_response.status_code}")
                                                        break  # 下载失败，跳出重试循环
                                                break

                                            elif state == "failed":
                                                logger.error(f"MinerU: Extraction failed - {file_result.get('err_msg')}")
                                                break

                                            elif state in ["running", "waiting-file", "pending", "converting"]:
                                                if attempt < max_retries - 1:
                                                    logger.info(f"MinerU: Processing... state={state} [Attempt {attempt + 1}/{max_retries}]")
                                                    time.sleep(retry_interval)
                                                else:
                                                    logger.warning(f"MinerU: Timeout after {max_retries * retry_interval} seconds")
                                                    break
                                    else:
                                        logger.error(f"MinerU: API error - {result_data.get('msg')}")
                                        break
                                else:
                                    logger.error(f"MinerU: Failed to get result, status={result_response.status_code}")
                                    break
                        else:
                            logger.error(f"MinerU: File upload failed, status={upload_file_response.status_code}")
                    else:
                        logger.error(f"MinerU: Failed to get upload URL - {upload_result.get('msg')}")
                else:
                    logger.error(f"MinerU: Failed to apply upload URL, status={upload_response.status_code}")

            except Exception as mineru_error:
                logger.error(f"MinerU extraction failed for {filename}: {mineru_error}")
        else:
            logger.info(f"MINERU_TOKEN not found in environment, skipping MinerU extraction")

        # 3. 尝试使用 pymupdf4llm
        try:
            import pymupdf4llm
            logger.info(f"Attempting to extract {filename} using pymupdf4llm...")
            
            # 使用pymupdf4llm提取PDF内容为Markdown格式
            md_text = pymupdf4llm.to_markdown(
                filepath,
                force_text=True,
                show_progress=False,  # 不显示进度条
                write_images=False,   # 不写出图片
                embed_images=False    # 不嵌入图片
            )
            
            if md_text:
                # 清理文本
                md_text = clean_text(md_text)
                logger.info(f"pymupdf4llm extracted: {md_text[:50]}...")  # 打印前 50 个字符
                
                if md_text and not self.is_text_garbled(md_text):
                    result['file_content'] = md_text
                    return result
                logger.info(f"pymupdf4llm result is empty or garbled for {filename}")
        except ImportError:
            logger.info("pymupdf4llm not installed. Skipping this extraction method.")
        except Exception as e:
            logger.info(f"pymupdf4llm failed for {filename}: {e}")

        # 4. 尝试使用 PyMuPDF (fitz)
        try:
            document = fitz.open(filepath)
            content = []
            for page in document:
                content.append(clean_text(page.get_text()))
            combined_text = clean_text("".join(content))
            logger.info(f"PyMuPDF extracted: {combined_text[:50]}...")  # 打印前 50 个字符
            if combined_text and not self.is_text_garbled(combined_text):
                result['file_content'] = combined_text
                return result
        except Exception as e:
            logger.info(f"PyMuPDF (fitz) failed for {filename}: {e}")

        # 5. 尝试使用 pdfplumber
        try:
            with pdfplumber.open(filepath) as pdf:
                content = []
                for page in pdf.pages:
                    content.append(clean_text(page.extract_text()))
            combined_text = clean_text("".join(content))
            logger.info(f"pdfplumber extracted: {combined_text[:50]}...")  # 打印前 50 个字符
            if combined_text and not self.is_text_garbled(combined_text):
                result['file_content'] = combined_text
                return result
        except Exception as e:
            logger.info(f"pdfplumber failed for {filename}: {e}")

        # 6. 尝试使用 pypdf
        try:
            with open(filepath, 'rb') as file:
                reader = pypdf.PdfReader(file)
                content = []
                for page in reader.pages:
                    content.append(clean_text(page.extract_text()))
            combined_text = clean_text("".join(content))
            logger.info(f"pypdf extracted: {combined_text[:50]}...")  # 打印前 50 个字符
            if combined_text and not self.is_text_garbled(combined_text):
                result['file_content'] = combined_text
                return result
        except Exception as e:
            logger.info(f"pypdf failed for {filename}: {e}")

        # 如果所有方法都失败或提取的内容为乱码，返回空
        logger.info(f"All extraction methods failed or text is garbled for {filename}")
        return {}

    def read_docx(self, filepath):
        try:
            filename = os.path.basename(filepath)
            result = {'file_extension': 'docx', 'file_name': filename}
            doc = docx.Document(filepath)
            full_text = []

            # Iterate through all elements in document body in order
            for element in doc.element.body:
                # Check if it's a paragraph
                if element.tag.endswith('p'):
                    para = docx.text.paragraph.Paragraph(element, doc)
                    if para.text.strip():
                        full_text.append(para.text)

                # Check if it's a table
                elif element.tag.endswith('tbl'):
                    table = docx.table.Table(element, doc)
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_text.append(cell_text)
                        if row_text:
                            full_text.append(' | '.join(row_text))

            result['file_content'] = '\n'.join(full_text)
            logger.info(f"result: {result}")
            return result
        except Exception as e:
            logger.error(f"""DOCX文件读取失败:
文件名: {filename if 'filename' in locals() else os.path.basename(filepath)}
文件路径: {filepath}
错误类型: {type(e).__name__}
错误信息: {str(e)}""")
            return {'file_extension': 'docx', 'file_name': filename, 'file_content': f'Error reading {filepath}: {e}'}


    def read_ppt(self, filepath):
        """读取PPT/PPTX文件内容"""
        try:
            filename = os.path.basename(filepath)
            file_extension = os.path.splitext(filepath)[1].lower()
            result = {'file_extension': file_extension[1:], 'file_name': filename}
            
            if not PPTX_AVAILABLE:
                logger.warning(f"python-pptx库不可用，无法处理PPT文件: {filename}")
                result['file_content'] = f"PPT文件: {filename} (无法处理，python-pptx库不可用)"
                return result
            
            # 只支持.pptx格式，.ppt格式需要转换
            if file_extension == '.ppt':
                logger.warning(f"不支持旧版PPT格式(.ppt)，请转换为.pptx格式: {filename}")
                result['file_content'] = f"PPT文件: {filename} (不支持.ppt格式，请转换为.pptx格式)"
                return result
            
            # 读取PPTX文件
            prs = Presentation(filepath)
            full_text = []
            
            # 提取每一页的内容
            for i, slide in enumerate(prs.slides, 1):
                slide_content = []
                slide_content.append(f"\n=== 第 {i} 页 ===")
                
                # 提取幻灯片中的所有文本
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                    
                    # 处理表格
                    if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                        try:
                            table = shape.table
                            table_content = []
                            for row in table.rows:
                                row_content = []
                                for cell in row.cells:
                                    if cell.text.strip():
                                        row_content.append(cell.text.strip())
                                if row_content:
                                    table_content.append(" | ".join(row_content))
                            if table_content:
                                slide_content.append("表格内容:")
                                slide_content.extend(table_content)
                        except Exception as e:
                            logger.debug(f"处理表格时出错: {e}")
                
                # 处理备注
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_content.append(f"备注: {notes_text}")
                
                if len(slide_content) > 1:  # 除了页码标题外还有其他内容
                    full_text.extend(slide_content)
            
            # 过滤空内容
            full_text = [item for item in full_text if item.strip()]
            
            if full_text:
                result['file_content'] = '\n'.join(full_text)
                logger.info(f"成功提取PPT内容: {filename}, 共 {len(prs.slides)} 页")
            else:
                result['file_content'] = f"PPT文件: {filename} (未提取到文本内容)"
                logger.warning(f"PPT文件未提取到内容: {filename}")
            
            return result
            
        except Exception as e:
            logger.error(f"""PPT文件读取失败:
文件名: {filename if 'filename' in locals() else os.path.basename(filepath)}
文件路径: {filepath}
错误类型: {type(e).__name__}
错误信息: {str(e)}""")
            filename = os.path.basename(filepath) if 'filename' not in locals() else filename
            file_extension = os.path.splitext(filepath)[1].lower()[1:] if 'file_extension' not in locals() else file_extension[1:]
            return {
                'file_extension': file_extension,
                'file_name': filename,
                'file_content': f'PPT文件读取失败: {str(e)}'
            }

    def read_txt(self, filepath):
        try:
            filename = os.path.basename(filepath)
            file_extension = os.path.splitext(filename)[1].lower()
            result = {'file_extension': file_extension, 'file_name': filename}

            with open(filepath, 'r', encoding='utf-8') as file:
                file_content = file.read()
            result['file_content'] = file_content
            return result
        except Exception as e:
            logger.error(f"""TXT文件读取失败:
文件名: {filename if 'filename' in locals() else os.path.basename(filepath)}
文件路径: {filepath}
错误类型: {type(e).__name__}
错误信息: {str(e)}""")
            return {'file_extension': '', 'file_name': '', 'file_content': ''}
